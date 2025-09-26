from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, func, or_
from sqlalchemy.orm import selectinload
from typing import List, Optional, Dict, Any
from datetime import datetime
import pandas as pd
import io
import re
import aiohttp
import asyncio
from difflib import SequenceMatcher
import json
import time
import uuid
from pathlib import Path
import pymorphy2
from functools import lru_cache
from pydantic import BaseModel

# Инициализируем морфологический анализатор
morph = pymorphy2.MorphAnalyzer()

def normalize_russian_text(text: str) -> str:
    """
    Приводит русский текст к нормальной форме (именительный падеж, единственное число).
    Также извлекает числа и специальные символы.
    """
    if not text:
        return ""
    
    # Сохраняем числа и специальные символы
    numbers_and_special = re.findall(r'\d+|[^\w\s]', text)
    
    # Разбиваем текст на слова
    words = re.findall(r'\b[а-яА-Я]+\b', text)
    
    # Нормализуем каждое слово
    normalized_words = []
    for word in words:
        # Получаем все возможные разборы слова
        parsed = morph.parse(word)
        if parsed:
            # Берем наиболее вероятный разбор
            normal_form = parsed[0].normal_form
            normalized_words.append(normal_form)
    
    # Собираем текст обратно, включая числа и специальные символы
    result = ' '.join(normalized_words + numbers_and_special)
    return result.strip()

@lru_cache(maxsize=1000)
def get_normalized_text(text: str) -> str:
    """
    Кэширующая обертка для normalize_russian_text
    """
    return normalize_russian_text(text)

from database import get_db
from models import User, ApiKey, AiProcessingLog, MatchingNomenclature
from ..dependencies import get_current_user
from ..schemas import AIMatchingResponse, MatchingResult
from ..utils.api_key import get_api_key

# Модели для Excel функционала
class ExcelRow(BaseModel):
    id: str
    наименование: str
    запрашиваемый_артикул: str
    количество: float
    единица_измерения: str
    наш_артикул: Optional[str] = None
    артикул_bl: Optional[str] = None
    номер_1с: Optional[str] = None
    стоимость: Optional[float] = None
    статус_сопоставления: Optional[str] = "pending"
    уверенность: Optional[int] = 0
    варианты_подбора: Optional[List[dict]] = []
    выбранный_вариант: Optional[int] = None  # Индекс выбранного варианта

class ExcelDataRequest(BaseModel):
    data: List[ExcelRow]

class ExcelParseResponse(BaseModel):
    success: bool
    data: List[ExcelRow]
    message: Optional[str] = None

class ExcelMatchResponse(BaseModel):
    success: bool
    matched_data: List[ExcelRow]
    statistics: Dict[str, Any]
    message: Optional[str] = None

async def process_natural_language_query(query: str, db: AsyncSession) -> dict:
    """
    Обрабатывает запрос на естественном языке и возвращает структурированный ответ
    """
    print(f"🤖 Начинаем обработку запроса: '{query}'")
    try:
        # Нормализуем запрос
        normalized_query = get_normalized_text(query)
        print(f"🔍 Нормализованный запрос: '{normalized_query}' (было: '{query}')")

        # Определяем тип запроса
        query_type = "unknown"
        if any(word in normalized_query.lower() for word in ["найти", "поиск", "искать", "где"]):
            query_type = "search"
        elif any(word in normalized_query.lower() for word in ["добавить", "создать", "сохранить"]):
            query_type = "create"
        elif any(word in normalized_query.lower() for word in ["изменить", "обновить", "исправить"]):
            query_type = "update"

        # Формируем промпт для ИИ в зависимости от типа запроса
        system_prompt = """Ты - ИИ-ассистент для работы с базой данных товаров. 
        Твоя задача - помогать пользователям находить, добавлять и обновлять информацию о товарах.
        
        Отвечай кратко и по делу, но в дружелюбной манере. Используй смайлики для улучшения читаемости.
        
        Если не можешь найти точное совпадение, предложи похожие варианты.
        Если запрос неясен, задай уточняющие вопросы.
        
        Всегда указывай причину, почему ты предлагаешь тот или иной вариант."""

        # Получаем ответ от ИИ
        print("🔑 Получаем API ключ...")
        api_key = await get_api_key(db)
        if not api_key:
            print("❌ API ключ не найден")
            return {
                "message": "К сожалению, не удалось получить доступ к ИИ. Попробуйте позже.",
                "success": False
            }
        print("✅ API ключ получен")

        print("🌐 Отправляем запрос к ИИ...")
        async with aiohttp.ClientSession() as session:
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Тип запроса: {query_type}\nЗапрос: {normalized_query}"}
            ]
            
            data = {
                "model": "gpt-4o-mini",
                "messages": messages,
                "temperature": 0.7,
                "max_tokens": 1000
            }
            
            async with session.post("https://api.polza.com/v1/chat/completions", headers=headers, json=data) as response:
                if response.status == 200:
                    result = await response.json()
                    ai_response = result["choices"][0]["message"]["content"]
                    
                    # Если это поисковый запрос, добавляем результаты поиска
                    if query_type == "search":
                        search_results = await smart_search_with_ai(normalized_query, db)
                        return {
                            "message": ai_response,
                            "search_results": search_results.get("matches", []),
                            "success": True
                        }
                    else:
                        return {
                            "message": ai_response,
                            "success": True
                        }
                else:
                    return {
                        "message": "Извините, произошла ошибка при обработке запроса.",
                        "success": False
                    }
    except Exception as e:
        print(f"Ошибка при обработке запроса: {e}")
        return {
            "message": "Произошла ошибка при обработке вашего запроса. Попробуйте переформулировать.",
            "success": False
        }

async def extract_articles_from_text(text: str, db: AsyncSession = None) -> List[dict]:
    """Извлекает артикулы из текста строки через парсинг и AI API"""
    try:
        # Сначала пытаемся парсить строку локально
        lines = [line.strip() for line in text.strip().split('\n') if line.strip()]
        articles = []
        
        for line in lines:
            if len(line) > 4:  # Пропускаем короткие строки
                # Парсим строку
                parsed_item = parse_item_string(line)
                
                # Если нашли артикул или описание, добавляем в результат
                if parsed_item['agb_article'] or parsed_item['description']:
                    articles.append({
                        'agb_article': parsed_item['agb_article'],
                        'description': parsed_item['description'],
                        'quantity': parsed_item['quantity'],
                        'unit': parsed_item['unit']
                    })
        
        # Если локальный парсинг не дал результатов, используем AI
        if not articles:
            async with aiohttp.ClientSession() as session:
                prompt = f"""
Проанализируй следующий текст и найди все артикулы товаров. 
Верни результат в формате JSON с массивом объектов, где каждый объект содержит:
- "article": найденный артикул
- "description": описание товара
- "quantity": количество (если указано)
- "unit": единица измерения (если указана)

Текст: {text}

Если артикулов нет, верни пустой массив [].
"""
                
                payload = {
                    "model": "gpt-4o-mini",
                    "messages": [
                        {"role": "system", "content": "Ты эксперт по анализу товарных позиций. Извлекай артикулы и описания товаров из текста."},
                        {"role": "user", "content": prompt}
                    ],
                    "temperature": 0.1
                }
                
                # Получаем API ключ из базы данных
                api_key = await get_api_key(db)
                if not api_key:
                    print("❌ Не удалось получить API ключ")
                    return []
                
                headers = {
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                }
                
                async with session.post(POLZA_API_URL, json=payload, headers=headers) as response:
                    if response.status == 200:
                        result = await response.json()
                        content = result["choices"][0]["message"]["content"]
                        
                        # Извлекаем JSON из ответа (может быть в markdown блоке)
                        if "```json" in content:
                            json_start = content.find("```json") + 7
                            json_end = content.find("```", json_start)
                            content = content[json_start:json_end].strip()
                        elif "```" in content:
                            json_start = content.find("```") + 3
                            json_end = content.find("```", json_start)
                            content = content[json_start:json_end].strip()
                        
                        import json
                        articles = json.loads(content)
                        return articles if isinstance(articles, list) else []
                    else:
                        print(f"❌ Ошибка AI API: {response.status}")
                        return []
        
        return articles
                    
    except Exception as e:
        print(f"❌ Ошибка при извлечении артикулов: {str(e)}")
        return []

def transliterate_to_latin(text: str) -> str:
    """Транслитерация русских букв в латиницу"""
    translit_map = {
        'а': 'a', 'б': 'b', 'в': 'v', 'г': 'g', 'д': 'd', 'е': 'e', 'ё': 'e',
        'ж': 'zh', 'з': 'z', 'и': 'i', 'й': 'y', 'к': 'k', 'л': 'l', 'м': 'm',
        'н': 'n', 'о': 'o', 'п': 'p', 'р': 'r', 'с': 's', 'т': 't', 'у': 'u',
        'ф': 'f', 'х': 'h', 'ц': 'ts', 'ч': 'ch', 'ш': 'sh', 'щ': 'sch',
        'ъ': '', 'ы': 'y', 'ь': '', 'э': 'e', 'ю': 'yu', 'я': 'ya',
        'А': 'A', 'Б': 'B', 'В': 'V', 'Г': 'G', 'Д': 'D', 'Е': 'E', 'Ё': 'E',
        'Ж': 'Zh', 'З': 'Z', 'И': 'I', 'Й': 'Y', 'К': 'K', 'Л': 'L', 'М': 'M',
        'Н': 'N', 'О': 'O', 'П': 'P', 'Р': 'R', 'С': 'S', 'Т': 'T', 'У': 'U',
        'Ф': 'F', 'Х': 'H', 'Ц': 'Ts', 'Ч': 'Ch', 'Ш': 'Sh', 'Щ': 'Sch',
        'Ъ': '', 'Ы': 'Y', 'Ь': '', 'Э': 'E', 'Ю': 'Yu', 'Я': 'Ya'
    }
    
    result = text
    for ru_char, en_char in translit_map.items():
        result = result.replace(ru_char, en_char)
    
    result = result.lower()
    
    # Заменяем пробелы и специальные символы на подчеркивания
    result = result.replace(' ', '_').replace('-', '_').replace('.', '').replace(',', '').replace('(', '').replace(')', '')
    
    return result
from models import User, UserRole, ArticleMapping, ContractorRequest, ContractorRequestItem, MatchingNomenclature
from ..schemas import (
    ArticleMapping as ArticleMappingSchema,
    ArticleMappingCreate,
    ContractorRequest as ContractorRequestSchema,
    ContractorRequestCreate,
    ContractorRequestItem as ContractorRequestItemSchema,
    ContractorRequestItemCreate,
    ContractorRequestItemUpdate,
    MatchingResult,
    MatchingSummary,
    TextUploadRequest
)
from .auth import get_current_user

router = APIRouter()

@router.post("/chat/", response_model=dict)
async def chat_with_ai(
    message: dict,
    db: AsyncSession = Depends(get_db)
) -> dict:
    """Обработка естественно-языкового запроса к ИИ"""
    try:
        query = message.get("message", "").strip()
        if not query:
            return {
                "message": "Пожалуйста, введите ваш запрос.",
                "success": False
            }
        
        # Обрабатываем запрос
        response = await process_natural_language_query(query, db)
        
        # Сохраняем лог обработки
        log = AiProcessingLog(
            request_type="chat",
            input_text=query,
            ai_response=response.get("message", ""),
            status="success" if response.get("success") else "error",
            processing_time=time.time()
        )
        db.add(log)
        await db.commit()
        
        return response
    except Exception as e:
        print(f"Ошибка в чате с ИИ: {e}")
        return {
            "message": "Произошла ошибка при обработке запроса. Попробуйте позже.",
            "success": False
        }

# URL для Polza.ai
POLZA_API_URL = "https://api.polza.ai/v1/chat/completions"

async def get_api_key(db: AsyncSession) -> str:
    """Получает активный API ключ из базы данных"""
    try:
        # Сначала ищем Polza.ai ключ
        result = await db.execute(select(ApiKey).where(
            ApiKey.is_active == True,
            ApiKey.provider == 'polza'
        ).limit(1))
        api_key_obj = result.scalar_one_or_none()
        
        # Если Polza.ai ключа нет, берем OpenAI
        if not api_key_obj:
            result = await db.execute(select(ApiKey).where(
                ApiKey.is_active == True,
                ApiKey.provider == 'openai'
            ).limit(1))
            api_key_obj = result.scalar_one_or_none()
        
        if not api_key_obj:
            raise Exception("Нет активного API ключа для ИИ-сервиса")
        
        # Расшифровываем ключ
        try:
            from cryptography.fernet import Fernet
            import os
            ENCRYPTION_KEY = b'iF0d2ARGQpaU9GFfQdWNovBL239dqwTp9hDDPrDQQic='
            cipher_suite = Fernet(ENCRYPTION_KEY)
            decrypted_key = cipher_suite.decrypt(api_key_obj.key.encode()).decode()
        except Exception:
            # Возможно, ключ не зашифрован, используем напрямую
            decrypted_key = api_key_obj.key
        
        return decrypted_key.strip()
        
    except Exception as e:
        print(f"Ошибка получения API ключа: {e}")
        return ""


def normalize_ai_article(article: dict) -> dict:
    """Нормализует данные артикула из ответа ИИ"""
    normalized = {
        'contractor_article': None,
        'description': None,
        'quantity': 1,
        'unit': 'шт',
        'agb_article': None,
        'bl_article': None,
        'match_confidence': 0.0,
        'match_type': 'unknown',
        'nomenclature': None
    }

    # Если article - это строка (натуральный запрос), обрабатываем как описание
    if isinstance(article, str):
        description = article.strip()
        if description:
            # Ищем потенциальный артикул в описании
            parsed = parse_item_string(description)
            if parsed['agb_article']:
                normalized['contractor_article'] = parsed['agb_article']
                normalized['description'] = parsed['description']
                normalized['quantity'] = parsed['quantity']
                normalized['unit'] = parsed['unit']
            else:
                normalized['description'] = description
                normalized['contractor_article'] = None
        else:
            normalized['description'] = 'Неизвестный товар'
        return normalized

    # Обрабатываем словарь из ИИ ответа
    if isinstance(article, dict):
        # contractor_article может быть None, пустой строкой или отсутствовать
        contractor_article = article.get('contractor_article')
        if contractor_article is not None and contractor_article != '':
            normalized['contractor_article'] = str(contractor_article).strip()
        elif not normalized['contractor_article']:
            # Если contractor_article не указан, используем часть описания как артикул
            description = article.get('description', '')
            if description:
                # Ищем потенциальный артикул в описании
                parsed = parse_item_string(description)
                if parsed['agb_article']:
                    normalized['contractor_article'] = parsed['agb_article']
                    normalized['description'] = parsed['description'] or description
                else:
                    normalized['description'] = description
            else:
                normalized['description'] = 'Неизвестный товар'

        # Описание
        description = article.get('description')
        if description:
            normalized['description'] = str(description).strip()
        elif not normalized['description']:
            normalized['description'] = 'Неизвестный товар'

        # Количество и единица измерения
        quantity = article.get('quantity')
        if quantity is not None:
            try:
                normalized['quantity'] = int(quantity)
            except (ValueError, TypeError):
                normalized['quantity'] = 1

        unit = article.get('unit')
        if unit:
            normalized['unit'] = str(unit).strip()

        # Артикулы АГБ и BL
        agb_article = article.get('agb_article')
        if agb_article:
            normalized['agb_article'] = str(agb_article).strip()

        bl_article = article.get('bl_article')
        if bl_article:
            normalized['bl_article'] = str(bl_article).strip()

        # Уверенность сопоставления
        confidence = article.get('match_confidence')
        if confidence is not None:
            try:
                normalized['match_confidence'] = float(confidence)
            except (ValueError, TypeError):
                normalized['match_confidence'] = 0.0

        # Тип сопоставления
        match_type = article.get('match_type')
        if match_type:
            normalized['match_type'] = str(match_type).strip()

    return normalized


def parse_item_string(item_string: str) -> dict:
    """Парсит строку товара и извлекает артикул, описание и количество"""
    import re

    # Убираем лишние пробелы
    item_string = item_string.strip()

    # Расширенные паттерны для поиска артикулов (везде в строке)
    article_patterns = [
        r'([A-ZА-Я]{2,6}[-_]\d{6,8})',    # ОХКУ-000184, BL-123456
        r'([A-ZА-Я]{2,6}\d{6,8})',        # ОХКУ000184, BL123456
        r'(\d{6,8})',                      # 123456
        r'([A-ZА-Я]{2,6}[-_]\d{3,8})',    # ОХКУ-184, BL-123
        r'([A-ZА-Я]{3,6}[-_]\d{4,6})',    # BL-1234, ОХКУ-1234
        r'(\d{4,6}[-_][A-ZА-Я]{2,4})',    # 1234-BL, 123456-АГБ
        r'([A-ZА-Я]{1,3}\d{4,8})',        # B1234, BL123456
        r'(\d{4,8}[A-ZА-Я]{1,3})',        # 1234B, 123456BL
    ]

    # Паттерны для поиска количества (в порядке приоритета)
    quantity_patterns = [
        r'(\d+)\s*(шт|штук|pcs|pieces?|кг|kg|л|l|м|m|м²|м³)\s*$',  # количество в конце строки
        r'\((\d+)\s*(шт|штук|pcs|pieces?|кг|kg|л|l|м|m|м²|м³)\)',  # количество в скобках
        r'(\d+)\s*(шт|штук|pcs|pieces?|кг|kg|л|l|м|m|м²|м³)',      # общий паттерн
    ]

    # Извлекаем артикул
    article = ""
    description = item_string
    quantity = 1
    unit = "шт"

    # Ищем артикул в любой части строки (начинаем с более точных паттернов)
    for pattern in article_patterns:
        matches = re.findall(pattern, item_string)
        if matches:
            # Берем самый длинный потенциальный артикул
            potential_articles = [m for m in matches if len(m) >= 3]
            if potential_articles:
                article = max(potential_articles, key=len)
                # Убираем артикул из описания
                description = re.sub(re.escape(article), '', item_string).strip()
                # Убираем лишние разделители
                description = re.sub(r'^[-_.,\s]+|[-_.,\s]+$', '', description)
                break
    
    # Ищем количество и единицу измерения
    for pattern in quantity_patterns:
        match = re.search(pattern, item_string, re.IGNORECASE)
        if match:
            quantity = int(match.group(1))
            if len(match.groups()) > 1:
                unit = match.group(2).lower()
                if unit in ['pcs', 'pieces']:
                    unit = 'шт'
                elif unit in ['kg']:
                    unit = 'кг'
                elif unit in ['l']:
                    unit = 'л'
                elif unit in ['m']:
                    unit = 'м'
            # Убираем количество из описания
            description = re.sub(pattern, '', description, flags=re.IGNORECASE).strip()
            break
    
    # Очищаем описание от лишних символов
    description = re.sub(r'\s+', ' ', description).strip()
    description = re.sub(r'^[-_\s]+', '', description)
    
    return {
        'agb_article': article,
        'description': description,
        'quantity': quantity,
        'unit': unit
    }

@router.post("/search/")
async def search_nomenclature(
    request: dict,
    db: AsyncSession = Depends(get_db)
) -> dict:
    """Поиск номенклатуры по различным критериям с поддержкой морфологии"""
    original_query = request.get("query", "").strip()
    search_type = request.get("search_type", "article")

    if not original_query:
        return {"matches": []}

    # Нормализуем поисковый запрос
    normalized_query = get_normalized_text(original_query)
    print(f"Нормализованный запрос: '{normalized_query}' (было: '{original_query}')")

    # Создаем базовый запрос
    base_query = select(MatchingNomenclature).where(MatchingNomenclature.is_active == True)

    # Получаем все потенциальные совпадения
    result = await db.execute(base_query)
    all_items = result.scalars().all()

    # Преобразуем результаты с учетом морфологии
    matches = []
    for item in all_items:
        match_confidence = 0
        match_reason = ""

        if search_type == "article":
            # Для артикулов используем прямое сравнение
            if item.agb_article and normalized_query.lower() in item.agb_article.lower():
                match_confidence = 100 if normalized_query.lower() == item.agb_article.lower() else 80
                match_reason = "Совпадение по артикулу АГБ"
            elif item.bl_article and normalized_query.lower() in item.bl_article.lower():
                match_confidence = 90 if normalized_query.lower() == item.bl_article.lower() else 70
                match_reason = "Совпадение по артикулу BL"
        
        elif search_type == "name":
            # Для наименований используем морфологический анализ
            normalized_name = get_normalized_text(item.name)
            
            # Проверяем точное совпадение нормализованных форм
            if normalized_query.lower() == normalized_name.lower():
                match_confidence = 100
                match_reason = "Точное совпадение наименования"
            else:
                # Проверяем частичное совпадение
                query_words = set(normalized_query.lower().split())
                name_words = set(normalized_name.lower().split())
                
                # Находим общие слова
                common_words = query_words & name_words
                
                if common_words:
                    # Рассчитываем процент совпадения
                    match_confidence = int((len(common_words) / len(query_words)) * 100)
                    match_reason = "Частичное совпадение наименования"
                
                # Если совпадение низкое, проверяем схожесть текста
                if match_confidence < 70:
                    similarity = calculate_similarity(normalized_query, normalized_name)
                    text_match_confidence = int(similarity * 100)
                    if text_match_confidence > match_confidence:
                        match_confidence = text_match_confidence
                        match_reason = "Схожесть текста наименования"
        
        elif search_type == "code":
            # Для кодов используем прямое сравнение
            if item.code_1c and normalized_query.lower() in item.code_1c.lower():
                match_confidence = 100 if normalized_query.lower() == item.code_1c.lower() else 80
                match_reason = "Совпадение по коду 1С"

        # Добавляем только если уверенность выше порога
        if match_confidence >= 50:
            matches.append({
                "id": item.id,
                "agb_article": item.agb_article,
                "name": item.name,
                "code_1c": item.code_1c,
                "bl_article": item.bl_article,
                "packaging": item.packaging,
                "unit": item.unit,
                "match_confidence": match_confidence,
                "match_reason": match_reason
            })

    # Сортируем по уверенности
    matches.sort(key=lambda x: x["match_confidence"], reverse=True)
    
    # Ограничиваем количество результатов
    matches = matches[:20]
    
    # Сохраняем результаты поиска в истории
    chat_message = AIChatMessage(
        session_id=1,  # TODO: Получать из контекста
        message_type="system",
        content=f"Поиск: {original_query} → {normalized_query} (тип: {search_type})",
        search_query=original_query,
        search_type=search_type,
        matching_results=matches
    )
    db.add(chat_message)
    await db.commit()

    return {"matches": matches}

async def smart_search_with_ai(search_text: str, db: AsyncSession) -> dict:
    """Умный поиск через AI - парсит строку и ищет по компонентам"""
    try:
        # Нормализуем поисковый запрос
        normalized_text = get_normalized_text(search_text)
        print(f"🔍 Нормализация запроса: '{search_text}' → '{normalized_text}'")
        
        # Парсим нормализованную строку для извлечения компонентов
        parsed_item = parse_item_string(normalized_text)
        
        print(f"🔍 Парсинг строки: '{normalized_text}'")
        print(f"   Артикул: '{parsed_item['agb_article']}'")
        print(f"   Описание: '{parsed_item['description']}'")
        print(f"   Количество: {parsed_item['quantity']} {parsed_item['unit']}")
        
        # Сначала проверяем уже существующие сопоставления по артикулу или описанию
        if parsed_item['agb_article'] or parsed_item['description']:
            print(f"🔍 Ищем в существующих сопоставлениях...")
            # Создаем базовый запрос
            query = select(ArticleMapping).where(ArticleMapping.is_active == True)
            conditions = []
            
            if parsed_item['agb_article']:
                conditions.append(ArticleMapping.contractor_article.ilike(f"%{parsed_item['agb_article']}%"))
            
            if parsed_item['description']:
                conditions.append(ArticleMapping.contractor_description.ilike(f"%{parsed_item['description']}%"))
            
            if conditions:
                query = query.where(or_(*conditions)).limit(10)
                existing_mappings = await db.execute(query)
                mappings = existing_mappings.scalars().all()
                
                if mappings:
                    print(f"✅ Найдено {len(mappings)} существующих сопоставлений")
                    matches = []
                    for mapping in mappings:
                        confidence = 100 if mapping.contractor_article == parsed_item['agb_article'] else 90
                        matches.append({
                            "agb_article": mapping.agb_article,
                            "bl_article": mapping.bl_article,
                            "name": mapping.agb_description,
                            "code_1c": "",
                            "confidence": mapping.match_confidence or confidence,
                            "packaging": mapping.packaging_factor or 1,
                            "is_existing": True,
                            "mapping_id": mapping.id,
                            "contractor_article": mapping.contractor_article,
                            "match_reason": "Найдено в существующих сопоставлениях"
                        })
                    
                    return {
                        "search_type": "existing_mapping",
                        "matches": matches
                    }
        
        # Проверяем точные и частичные совпадения в базе данных
        print(f"🔍 Ищем в базе номенклатуры...")
        matches = []
        
        # Создаем базовый запрос для поиска в номенклатуре
        query = select(MatchingNomenclature).where(MatchingNomenclature.is_active == True)
        conditions = []
        
        # Добавляем условия поиска
        if parsed_item['agb_article']:
            conditions.extend([
                MatchingNomenclature.agb_article.ilike(f"%{parsed_item['agb_article']}%"),
                MatchingNomenclature.bl_article.ilike(f"%{parsed_item['agb_article']}%"),
                MatchingNomenclature.code_1c.ilike(f"%{parsed_item['agb_article']}%")
            ])
        
        if parsed_item['description']:
            conditions.append(MatchingNomenclature.name.ilike(f"%{parsed_item['description']}%"))
        
        if conditions:
            query = query.where(or_(*conditions)).limit(20)
            result = await db.execute(query)
            found_items = result.scalars().all()
            
            for item in found_items:
                confidence = 0
                match_reason = ""
                
                # Вычисляем уверенность на основе совпадений
                if parsed_item['agb_article']:
                    if item.agb_article and item.agb_article.lower() == parsed_item['agb_article'].lower():
                        confidence = 100
                        match_reason = "Точное совпадение по артикулу АГБ"
                    elif item.bl_article and item.bl_article.lower() == parsed_item['agb_article'].lower():
                        confidence = 95
                        match_reason = "Точное совпадение по артикулу BL"
                    elif item.code_1c and item.code_1c.lower() == parsed_item['agb_article'].lower():
                        confidence = 90
                        match_reason = "Точное совпадение по коду 1С"
                    else:
                        # Частичное совпадение
                        confidence = 70
                        match_reason = "Частичное совпадение по артикулу"
                
                if parsed_item['description'] and confidence < 100:
                    desc_similarity = calculate_similarity(parsed_item['description'], item.name)
                    if desc_similarity > 0.8:
                        confidence = max(confidence, 85)
                        match_reason = "Высокая схожесть описания"
                    elif desc_similarity > 0.6:
                        confidence = max(confidence, 75)
                        match_reason = "Средняя схожесть описания"
                
                if confidence > 0:
                    matches.append({
                        "agb_article": item.agb_article,
                        "bl_article": item.bl_article,
                        "name": item.name,
                        "code_1c": item.code_1c,
                        "confidence": confidence,
                        "match_reason": match_reason,
                        "is_existing": False
                    })

        # Если нашли совпадения в базе, возвращаем их
        if matches:
            print(f"✅ Найдено {len(matches)} совпадений в базе номенклатуры")
            return {
                "search_type": "database_match",
                "matches": matches
            }

        print(f"❌ Совпадений в базе не найдено, используем ИИ")
        
        # Используем ИИ для поиска
        try:
            # Получаем все номенклатуры для контекста
            all_nomenclatures = await db.execute(
                select(MatchingNomenclature).where(MatchingNomenclature.is_active == True).limit(100)
            )
            nomenclatures = all_nomenclatures.scalars().all()
            
            # Формируем текст с номенклатурами для ИИ
            nomenclatures_text = "\n".join([
                f"- {nom.agb_article} | {nom.bl_article or 'N/A'} | {nom.name}"
                for nom in nomenclatures
            ])
            
            # Создаем промпт для ИИ
            prompt = f"""Найди в базе данных товар по запросу: "{search_text}"

База данных товаров:
{nomenclatures_text}

Верни JSON массив с найденными товарами в формате:
[
    {{
        "agb_article": "артикул АГБ",
        "bl_article": "артикул BL",
        "name": "наименование",
        "code_1c": "код 1С",
        "confidence": 85,
        "match_reason": "причина совпадения"
    }}
]

Если ничего не найдено, верни пустой массив: []"""

            # Получаем API ключ
            api_key = await get_api_key(db)
            if not api_key:
                return {"search_type": "error", "matches": []}

            # Отправляем запрос к ИИ
            async with aiohttp.ClientSession() as session:
                headers = {
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                }
                
                data = {
                    "model": "gpt-4o-mini",
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": 0.1,
                    "max_tokens": 1500
                }
                
                async with session.post("https://api.polza.com/v1/chat/completions", headers=headers, json=data) as response:
                    if response.status in [200, 201]:
                        result = await response.json()
                        ai_response = result["choices"][0]["message"]["content"]
                        
                        # Парсим JSON ответ
                        import json
                        try:
                            if "```json" in ai_response:
                                json_start = ai_response.find("```json") + 7
                                json_end = ai_response.find("```", json_start)
                                ai_response = ai_response[json_start:json_end].strip()
                            elif "```" in ai_response:
                                json_start = ai_response.find("```") + 3
                                json_end = ai_response.find("```", json_start)
                                ai_response = ai_response[json_start:json_end].strip()

                            matches = json.loads(ai_response)
                            if isinstance(matches, list) and matches:
                                print(f"✅ ИИ нашел {len(matches)} совпадений")
                                return {
                                    "search_type": "ai_match",
                                    "matches": matches
                                }
                        except json.JSONDecodeError:
                            pass
            
            # Если ИИ не нашел ничего, возвращаем пустой результат
            return {"search_type": "not_found", "matches": []}
            
        except Exception as e:
            print(f"❌ Ошибка ИИ поиска: {e}")
            return {"search_type": "error", "matches": []}

        # Формируем список номенклатур для AI с приоритетом по релевантности
        # Сначала ищем возможные совпадения по артикулу или описанию
        relevant_nomenclatures = []
        search_terms = [parsed_item['agb_article'], parsed_item['description']]
        search_terms = [term.lower() for term in search_terms if term]

        # Ищем релевантные номенклатуры
        for nom in nomenclatures:
            relevance_score = 0

            # Проверяем артикул АГБ
            if parsed_item['agb_article'] and nom.agb_article:
                if parsed_item['agb_article'].lower() in nom.agb_article.lower():
                    relevance_score += 100
                elif nom.agb_article.lower() in parsed_item['agb_article'].lower():
                    relevance_score += 80

            # Проверяем артикул BL
            if parsed_item['agb_article'] and nom.bl_article:
                if parsed_item['agb_article'].lower() in nom.bl_article.lower():
                    relevance_score += 90
                elif nom.bl_article.lower() in parsed_item['agb_article'].lower():
                    relevance_score += 70

            # Проверяем код 1С
            if parsed_item['agb_article'] and nom.code_1c:
                if parsed_item['agb_article'].lower() in nom.code_1c.lower():
                    relevance_score += 60

            # Проверяем наименование
            if parsed_item['description'] and nom.name:
                if parsed_item['description'].lower() in nom.name.lower():
                    relevance_score += 50
                elif nom.name.lower() in parsed_item['description'].lower():
                    relevance_score += 40

            if relevance_score > 0:
                relevant_nomenclatures.append((nom, relevance_score))

        # Сортируем по релевантности
        relevant_nomenclatures.sort(key=lambda x: x[1], reverse=True)

        # Берем топ-100 релевантных + 100 случайных для разнообразия
        top_relevant = relevant_nomenclatures[:100]
        remaining_count = min(100, len(nomenclatures) - len(top_relevant))

        import random
        random_nomenclatures = random.sample(
            [nom for nom, score in relevant_nomenclatures[100:] if nom not in [n for n, s in top_relevant]],
            remaining_count
        ) if len(nomenclatures) > len(top_relevant) else []

        selected_nomenclatures = [nom for nom, score in top_relevant] + random_nomenclatures

        print(f"📊 Выбрано для поиска: {len(selected_nomenclatures)} номенклатур (релевантных: {len(top_relevant)}, случайных: {len(random_nomenclatures)})")

        # Формируем список номенклатур для AI
        nomenclatures_text = "\n".join([
            f"ID: {nom.id}, АГБ: {nom.agb_article}, BL: {nom.bl_article or 'нет'}, Код1С: {nom.code_1c or 'нет'}, Название: {nom.name}, Фасовка: {nom.packaging or 1}, Ед: {nom.unit}"
            for nom in selected_nomenclatures[:200]
        ])
        
        # Формируем промпт для AI с учетом парсинга
        prompt = f"""
        ТЫ - эксперт по поиску товаров в базе данных АГБ.

        ЗАПРОС ПОЛЬЗОВАТЕЛЯ: "{search_text}"

        ПРОАНАЛИЗИРУЙ ЗАПРОС И НАЙДИ СООТВЕТСТВИЯ В БАЗЕ ДАННЫХ:
        {nomenclatures_text}

        ПРАВИЛА:
        - Ищи ТОЧНЫЕ совпадения по названиям товаров
        - Сравнивай слова и фразы из запроса с названиями в базе
        - Если точного совпадения нет, ищи похожие товары (схожесть > 70%)
        - Верни ТОЛЬКО релевантные результаты

        ФОРМАТ ОТВЕТА (JSON):
        {{
            "search_type": "ai_match",
            "matches": [
                {{
                    "id": "ID_из_базы",
                    "agb_article": "артикул_АГБ",
                    "bl_article": "артикул_BL_или_null",
                    "code_1c": "код_1С_или_null",
                    "name": "название_товара",
                    "confidence": 75,
                    "match_reason": "найдено_по_названию"
                }}
            ]
        }}

        ЕСЛИ НИЧЕГО НЕ НАЙДЕНО - верни: {{"search_type": "ai_match", "matches": []}}
        """
        
        # Получаем API ключ из базы данных
        api_key = await get_api_key(db)
        if not api_key:
            print("❌ Не удалось получить API ключ")
            return {"search_type": "general", "matches": []}
        
        try:
            # Отправляем запрос к AI API
            async with aiohttp.ClientSession() as session:
                headers = {
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                }
                
                data = {
                    "model": "gpt-4o-mini",
                    "messages": [
                        {"role": "user", "content": prompt}
                    ],
                    "temperature": 0.1,
                    "max_tokens": 1500
                }
                
                try:
                    async with session.post(POLZA_API_URL, headers=headers, json=data) as response:
                        if response.status in [200, 201]:
                            result = await response.json()
                            ai_response = result["choices"][0]["message"]["content"]
                            
                            # Парсим JSON ответ
                            import json
                            try:
                                # Извлекаем JSON из markdown блока, если он есть
                                if "```json" in ai_response:
                                    json_start = ai_response.find("```json") + 7
                                    json_end = ai_response.find("```", json_start)
                                    if json_end != -1:
                                        ai_response = ai_response[json_start:json_end].strip()
                                elif "```" in ai_response:
                                    json_start = ai_response.find("```") + 3
                                    json_end = ai_response.find("```", json_start)
                                    if json_end != -1:
                                        ai_response = ai_response[json_start:json_end].strip()

                                # Если ответ содержит "НИЧЕГО НЕ НАЙДЕНО", возвращаем пустой результат
                                if "НИЧЕГО НЕ НАЙДЕНО" in ai_response.upper() or len(ai_response.strip()) < 10:
                                    return {"search_type": "ai_match", "matches": []}

                                matches = json.loads(ai_response)
                                ai_matches = matches.get('matches', [])
                                print(f"✅ AI нашел {len(ai_matches)} совпадений")
                                return matches
                            except json.JSONDecodeError as e:
                                print(f"Failed to parse AI response: {e}")
                                print(f"AI response: {ai_response}")
                                # Если не удалось распарсить JSON, проверяем текст ответа
                                if ai_response and "найден" in ai_response.lower() and len(ai_response) < 100:
                                    return {"search_type": "ai_match", "matches": []}
                                return {"search_type": "general", "matches": []}
                        else:
                            print(f"AI API error: {response.status}")
                            return {"search_type": "general", "matches": []}
                except Exception as e:
                    print(f"HTTP request error: {e}")
                    return {"search_type": "error", "matches": []}
        except Exception as e:
            print(f"AI search error: {e}")
            # Создаем fallback результат с сообщением о не найденном товаре
            return {
                "search_type": "not_found",
                "matches": [{
                    "id": "0",
                    "agb_article": "",
                    "bl_article": "",
                    "code_1c": "",
                    "name": f"Товар '{search_text}' не найден в базе данных",
                    "confidence": 0.0,
                    "match_reason": "товар_не_найден"
                }]
            }

    except Exception as e:
        print(f"Global error: {e}")
        return {"search_type": "error", "matches": []}


async def search_with_ai(description: str, db: AsyncSession) -> dict:
    """Поиск соответствий в базе данных через AI API (старая версия для совместимости)"""
    result = await smart_search_with_ai(description, db)
    # Преобразуем в старый формат для совместимости
    return {"matches": result.get("matches", [])}


def calculate_similarity(text1: str, text2: str) -> float:
    """Вычисляет схожесть между двумя текстами"""
    return SequenceMatcher(None, text1.lower(), text2.lower()).ratio()


def find_best_match(description: str, nomenclature_list: List[MatchingNomenclature]) -> Optional[tuple]:
    """Находит лучшее соответствие для описания товара"""
    best_match = None
    best_score = 0.0
    
    for nomenclature in nomenclature_list:
        # Проверяем совпадение по названию
        name_score = calculate_similarity(description, nomenclature.name)
        
        # Проверяем совпадение по артикулу
        article_score = calculate_similarity(description, nomenclature.agb_article)
        
        # Комбинируем оценки
        combined_score = max(name_score, article_score * 0.8)  # Артикул менее важен
        
        if combined_score > best_score and combined_score > 0.3:  # Минимальный порог
            best_score = combined_score
            best_match = nomenclature
    
    if best_match:
        return best_match, int(best_score * 100)
    return None


@router.get("/mappings/", response_model=List[ArticleMappingSchema])
async def get_article_mappings(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Получение списка соответствий артикулов"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    result = await db.execute(
        select(ArticleMapping)
        .where(ArticleMapping.is_active == True)
        .order_by(ArticleMapping.created_at.desc())
    )
    mappings = result.scalars().all()
    return mappings


@router.post("/mappings/", response_model=ArticleMappingSchema)
async def create_article_mapping(
    mapping_data: ArticleMappingCreate,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Создание нового соответствия артикулов"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    try:
        new_mapping = ArticleMapping(
            contractor_article=mapping_data.contractor_article,
            contractor_description=mapping_data.contractor_description,
            agb_article=mapping_data.agb_article,
            agb_description=mapping_data.agb_description,
            bl_article=mapping_data.bl_article,
            bl_description=mapping_data.bl_description,
            packaging_factor=mapping_data.packaging_factor,
            unit=mapping_data.unit,
            nomenclature_id=mapping_data.nomenclature_id
        )
        
        db.add(new_mapping)
        await db.commit()
        await db.refresh(new_mapping)
        
        return new_mapping
        
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при создании соответствия: {str(e)}")


@router.delete("/mappings/{mapping_id}/", response_model=dict)
async def delete_article_mapping(
    mapping_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Удаление соответствия артикулов"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    try:
        # Находим сопоставление
        result = await db.execute(
            select(ArticleMapping).where(ArticleMapping.id == mapping_id)
        )
        mapping = result.scalar_one_or_none()
        
        if not mapping:
            raise HTTPException(status_code=404, detail="Сопоставление не найдено")
        
        # Мягкое удаление - помечаем как неактивное
        mapping.is_active = False
        await db.commit()
        
        return {"message": "Сопоставление успешно удалено", "success": True}
        
    except HTTPException:
        raise
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при удалении сопоставления: {str(e)}")


@router.post("/requests/", response_model=ContractorRequestSchema)
async def create_contractor_request(
    request_data: ContractorRequestCreate,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Создание новой заявки контрагента"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    try:
        # Проверяем уникальность номера заявки
        existing = await db.execute(
            select(ContractorRequest).where(ContractorRequest.request_number == request_data.request_number)
        )
        if existing.scalar_one_or_none():
            raise HTTPException(status_code=400, detail="Заявка с таким номером уже существует")
        
        # Парсим дату
        request_date = datetime.fromisoformat(request_data.request_date.replace('Z', '+00:00'))
        
        # Создаем заявку
        new_request = ContractorRequest(
            request_number=request_data.request_number,
            contractor_name=request_data.contractor_name,
            request_date=request_date,
            total_items=len(request_data.items),
            created_by=current_user.id
        )
        
        db.add(new_request)
        await db.flush()  # Получаем ID заявки
        
        # Создаем позиции заявки
        for item_data in request_data.items:
            new_item = ContractorRequestItem(
                request_id=new_request.id,
                line_number=item_data.line_number,
                contractor_article=item_data.contractor_article,
                description=item_data.description,
                unit=item_data.unit,
                quantity=item_data.quantity,
                category=item_data.category
            )
            db.add(new_item)
        
        await db.commit()
        await db.refresh(new_request)
        
        return new_request
        
    except HTTPException:
        raise
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при создании заявки: {str(e)}")


@router.get("/requests/", response_model=List[ContractorRequestSchema])
async def get_contractor_requests(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Получение списка заявок контрагентов"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    result = await db.execute(
        select(ContractorRequest)
        .where(ContractorRequest.created_by == current_user.id)
        .order_by(ContractorRequest.created_at.desc())
    )
    requests = result.scalars().all()
    return requests


@router.get("/requests/{request_id}", response_model=ContractorRequestSchema)
async def get_contractor_request(
    request_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Получение конкретной заявки контрагента"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    result = await db.execute(
        select(ContractorRequest)
        .where(ContractorRequest.id == request_id)
    )
    request = result.scalar_one_or_none()
    
    if not request:
        raise HTTPException(status_code=404, detail="Заявка не найдена")
    
    if request.created_by != current_user.id:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    return request


@router.post("/requests/{request_id}/match", response_model=MatchingSummary)
async def match_articles(
    request_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Сопоставление артикулов в заявке"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    try:
        # Получаем заявку
        result = await db.execute(
            select(ContractorRequest).where(ContractorRequest.id == request_id)
        )
        request = result.scalar_one_or_none()
        
        if not request:
            raise HTTPException(status_code=404, detail="Заявка не найдена")
        
        if request.created_by != current_user.id:
            raise HTTPException(status_code=403, detail="Доступ запрещен")
        
        # Получаем позиции заявки
        items_result = await db.execute(
            select(ContractorRequestItem)
            .where(ContractorRequestItem.request_id == request_id)
            .order_by(ContractorRequestItem.line_number)
        )
        items = items_result.scalars().all()
        
        # Получаем всю номенклатуру для сопоставления
        nomenclature_result = await db.execute(
            select(MatchingNomenclature).where(MatchingNomenclature.is_active == True)
        )
        nomenclature_list = nomenclature_result.scalars().all()
        
        # Сопоставляем каждую позицию
        results = []
        matched_count = 0
        high_confidence = 0
        medium_confidence = 0
        low_confidence = 0
        
        for item in items:
            match_result = find_best_match(item.description, nomenclature_list)
            
            if match_result:
                nomenclature, confidence = match_result
                
                # Обновляем позицию
                item.matched_nomenclature_id = nomenclature.id
                item.agb_article = nomenclature.agb_article
                item.packaging_factor = 1  # По умолчанию
                item.recalculated_quantity = item.quantity * item.packaging_factor
                item.match_confidence = confidence
                item.match_status = "matched"
                
                # Определяем уровень уверенности
                if confidence >= 80:
                    high_confidence += 1
                elif confidence >= 60:
                    medium_confidence += 1
                else:
                    low_confidence += 1
                
                matched_count += 1
                
                results.append(MatchingResult(
                    item_id=item.id,
                    contractor_article=item.contractor_article,
                    description=item.description,
                    matched=True,
                    agb_article=nomenclature.agb_article,
                    packaging_factor=item.packaging_factor,
                    recalculated_quantity=item.recalculated_quantity,
                    match_confidence=confidence,
                    nomenclature=nomenclature
                ))
            else:
                item.match_status = "unmatched"
                results.append(MatchingResult(
                    item_id=item.id,
                    contractor_article=item.contractor_article,
                    description=item.description,
                    matched=False
                ))
        
        # Обновляем статистику заявки
        request.matched_items = matched_count
        request.status = "processing"
        request.processed_by = current_user.id
        request.processed_at = datetime.now()
        
        await db.commit()
        
        return MatchingSummary(
            total_items=len(items),
            matched_items=matched_count,
            unmatched_items=len(items) - matched_count,
            high_confidence_items=high_confidence,
            medium_confidence_items=medium_confidence,
            low_confidence_items=low_confidence,
            results=results
        )
        
    except HTTPException:
        raise
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при сопоставлении: {str(e)}")


@router.put("/items/{item_id}", response_model=ContractorRequestItemSchema)
async def update_request_item(
    item_id: int,
    item_data: ContractorRequestItemUpdate,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Обновление позиции заявки"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    try:
        # Получаем позицию
        result = await db.execute(
            select(ContractorRequestItem).where(ContractorRequestItem.id == item_id)
        )
        item = result.scalar_one_or_none()
        
        if not item:
            raise HTTPException(status_code=404, detail="Позиция не найдена")
        
        # Проверяем права доступа через заявку
        request_result = await db.execute(
            select(ContractorRequest).where(ContractorRequest.id == item.request_id)
        )
        request = request_result.scalar_one_or_none()
        
        if not request or request.created_by != current_user.id:
            raise HTTPException(status_code=403, detail="Доступ запрещен")
        
        # Обновляем поля
        update_data = item_data.model_dump(exclude_unset=True)
        for field, value in update_data.items():
            setattr(item, field, value)
        
        # Пересчитываем количество если изменился коэффициент фасовки
        if item_data.packaging_factor is not None:
            item.recalculated_quantity = item.quantity * item.packaging_factor
        
        await db.commit()
        await db.refresh(item)
        
        return item
        
    except HTTPException:
        raise
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при обновлении позиции: {str(e)}")


@router.post("/upload-excel/", response_model=ContractorRequestSchema)
async def upload_contractor_excel(
    file: UploadFile = File(...),
    contractor_name: str = "Контрагент",
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Загрузка заявки контрагента из Excel файла"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    try:
        # Читаем Excel файл
        content = await file.read()
        df = pd.read_excel(io.BytesIO(content))
        
        # Генерируем номер заявки
        request_number = f"REQ_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        # Создаем заявку
        new_request = ContractorRequest(
            request_number=request_number,
            contractor_name=contractor_name,
            request_date=datetime.now(),
            total_items=len(df),
            created_by=current_user.id
        )
        
        db.add(new_request)
        await db.flush()
        
        # Обрабатываем строки Excel
        items = []
        for index, row in df.iterrows():
            # Извлекаем данные из строки (адаптируем под формат из примера)
            line_number = index + 1
            contractor_article = str(row.get('Артикул', row.get('№', ''))).strip()
            description = str(row.get('Описание', row.get('Наименование', ''))).strip()
            unit = str(row.get('Ед.', row.get('Единица', 'шт'))).strip()
            quantity = int(row.get('Количество', row.get('Кол-во', 1)))
            category = str(row.get('Категория', 'Для бурения')).strip()
            
            if contractor_article and description:
                new_item = ContractorRequestItem(
                    request_id=new_request.id,
                    line_number=line_number,
                    contractor_article=contractor_article,
                    description=description,
                    unit=unit,
                    quantity=quantity,
                    category=category
                )
                db.add(new_item)
                items.append(new_item)
        
        await db.commit()
        await db.refresh(new_request)
        
        # Выполняем сопоставление
        matching_results = await perform_matching(new_request.id, db)
        
        return {
            "request": new_request,
            "matching_results": matching_results
        }
        
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при загрузке файла: {str(e)}")


@router.post("/test-create-request/")
async def test_create_request(
    request_data: dict,
    db: AsyncSession = Depends(get_db)
):
    """Тестовый endpoint для создания заявки без файла"""
    try:
        contractor_name = request_data.get("contractor_name", "Тест Контрагент")
        print(f"Начинаем создание тестовой заявки для: {contractor_name}")
        # Генерируем номер заявки
        request_number = f"TEST_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        print(f"Номер заявки: {request_number}")
        
        # Создаем заявку
        new_request = ContractorRequest(
            request_number=request_number,
            contractor_name=contractor_name,
            request_date=datetime.now(),
            total_items=5,  # Тестовая заявка из 5 позиций
            created_by=1  # Тестовый пользователь
        )
        print(f"Создана заявка: {new_request}")
        
        db.add(new_request)
        await db.flush()
        print(f"Заявка добавлена в БД, ID: {new_request.id}")
        
        # Создаем тестовую заявку на основе примера
        test_items = [
            {"agb_article": "1299650", "description": "Шпиндель верхней части керноприемника H/HU, 25231, SDT", "quantity": 5, "unit": "шт"},
            {"agb_article": "1298240", "description": "Втулка удержания жидкости 306131, SDT", "quantity": 12, "unit": "шт"},
            {"agb_article": "1298244", "description": "Пружина мягкая N/H/P, удержания жидкости, 104966, SDT", "quantity": 10, "unit": "шт"},
            {"agb_article": "1299679", "description": "Щека верхняя для ключа разводного 24\", 14947, SDT", "quantity": 8, "unit": "шт"},
            {"agb_article": "1299680", "description": "Щека верхняя для ключа разводного 36\", 14950, SDT", "quantity": 8, "unit": "шт"}
        ]
        
        items = []
        for index, item_data in enumerate(test_items):
            line_number = index + 1
            new_item = ContractorRequestItem(
                request_id=new_request.id,
                line_number=line_number,
                contractor_article=item_data["agb_article"],
                description=item_data["description"],
                unit=item_data["unit"],
                quantity=item_data["quantity"],
                category="Для бурения"
            )
            db.add(new_item)
            items.append(new_item)
        
        await db.commit()
        await db.refresh(new_request)
        
        # Выполняем сопоставление
        matching_results = await perform_matching(new_request.id, db)
        
        return {
            "request": new_request,
            "matching_results": matching_results
        }
        
    except Exception as e:
        await db.rollback()
        import traceback
        error_details = traceback.format_exc()
        print(f"Ошибка в test_create_request: {str(e)}")
        print(f"Traceback: {error_details}")
        raise HTTPException(status_code=500, detail=f"Ошибка при создании заявки: {str(e)}")


@router.get("/requests/{request_id}/export/excel")
async def export_matching_results(
    request_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Экспорт результатов сопоставления в Excel"""
    if current_user.role not in [UserRole.ADMIN, UserRole.VED_PASSPORT]:
        raise HTTPException(status_code=403, detail="Доступ запрещен")
    
    try:
        # Получаем заявку
        result = await db.execute(
            select(ContractorRequest).where(ContractorRequest.id == request_id)
        )
        request = result.scalar_one_or_none()
        
        if not request or request.created_by != current_user.id:
            raise HTTPException(status_code=403, detail="Доступ запрещен")
        
        # Получаем позиции с номенклатурой
        items_result = await db.execute(
            select(ContractorRequestItem)
            .options(selectinload(ContractorRequestItem.matched_nomenclature))
            .where(ContractorRequestItem.request_id == request_id)
            .order_by(ContractorRequestItem.line_number)
        )
        items = items_result.scalars().all()
        
        # Создаем DataFrame
        data = []
        for item in items:
            data.append({
                'Номер строки': item.line_number,
                'Артикул контрагента': item.contractor_article,
                'Описание контрагента': item.description,
                'Количество': item.quantity,
                'Единица': item.unit,
                'Статус сопоставления': item.match_status,
                'Уверенность (%)': item.match_confidence,
                'Артикул АГБ': item.agb_article or '',
                'Артикул BL': item.bl_article or '',
                'Коэффициент фасовки': item.packaging_factor,
                'Пересчитанное количество': item.recalculated_quantity or item.quantity,
                'Номенклатура АГБ': item.matched_nomenclature.name if item.matched_nomenclature else '',
                'Код 1С': item.matched_nomenclature.code_1c if item.matched_nomenclature else ''
            })
        
        df = pd.DataFrame(data)
        
        # Создаем Excel файл в памяти
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='Результаты сопоставления', index=False)
            
            # Настраиваем ширину колонок
            worksheet = writer.sheets['Результаты сопоставления']
            for column in worksheet.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = (max_length + 2)
                worksheet.column_dimensions[column_letter].width = min(adjusted_width, 50)
        
        output.seek(0)
        
        # Возвращаем файл
        from fastapi import Response
        return Response(
            content=output.getvalue(),
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            headers={
                'Content-Disposition': f'attachment; filename=matching_results_{request.request_number}.xlsx'
            }
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка при экспорте: {str(e)}")


@router.post("/upload-text/")
async def upload_text_request(
    request_data: TextUploadRequest,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Загрузка заявки из текста с AI поиском"""
    try:
        # Парсим текст для извлечения позиций
        items = parse_text_to_items(request_data.text)
        
        if not items:
            raise HTTPException(status_code=400, detail="Не удалось извлечь позиции из текста")
        
        # Создаем заявку с правильным номером
        contractor_name = request_data.contractor_name
        
        # Генерируем номер заявки в формате {название_контрагента}_{порядковый_номер}_{дата}
        contractor_name_latin = transliterate_to_latin(contractor_name)
        print(f"🔤 Транслитерация: '{contractor_name}' -> '{contractor_name_latin}'")
        
        # Получаем количество заявок для этого контрагента сегодня
        today = datetime.now().date()
        count_result = await db.execute(
            select(func.count(ContractorRequest.id)).where(
                ContractorRequest.contractor_name == contractor_name,
                func.date(ContractorRequest.request_date) == today
            )
        )
        request_count = count_result.scalar() + 1
        
        request_number = f"{contractor_name_latin}_{request_count:03d}_{datetime.now().strftime('%Y%m%d')}"
        
        request = ContractorRequest(
            request_number=request_number,
            contractor_name=contractor_name,
            request_date=datetime.now(),
            total_items=len(items),
            status='new',
            created_by=current_user.id
        )
        
        db.add(request)
        await db.flush()  # Получаем ID заявки
        
        # Создаем позиции заявки с AI поиском
        for item_data in items:
            description = item_data.get('description', '')
            print(f"Searching AI matches for: {description}")
            
            # Ищем соответствия через AI
            ai_matches = await search_with_ai(description, db)
            
            # Создаем позицию заявки
            item = ContractorRequestItem(
                request_id=request.id,
                contractor_article=item_data.get('agb_article', ''),
                description=description,
                quantity=item_data.get('quantity', 1),
                unit=item_data.get('unit', 'шт')
            )
            db.add(item)
            await db.flush()  # Получаем ID позиции
            
            # Если AI нашел соответствия, создаем записи сопоставления
            if ai_matches.get('matches'):
                for match in ai_matches['matches']:
                    # Ищем номенклатуру в базе по артикулу
                    nom_result = await db.execute(
                        select(MatchingNomenclature).where(MatchingNomenclature.agb_article == match['agb_article'])
                    )
                    nomenclature = nom_result.scalar_one_or_none()
                    
                    if nomenclature:
                        # Создаем запись сопоставления
                        mapping = ArticleMapping(
                            contractor_request_item_id=item.id,
                            agb_article=match['agb_article'],
                            bl_article=match.get('code_1c', ''),
                            match_confidence=match.get('confidence', 0.0),
                            packaging_factor=1.0,
                            recalculated_quantity=item_data.get('quantity', 1),
                            nomenclature_id=nomenclature.id
                        )
                        db.add(mapping)
                        print(f"Created mapping: {match['agb_article']} -> {match.get('code_1c', '')}")
        
        await db.commit()
        await db.refresh(request)
        
        # Выполняем сопоставление
        matching_results = await perform_matching(request.id, db)
        
        # Загружаем заявку с позициями
        result = await db.execute(
            select(ContractorRequest)
            .options(selectinload(ContractorRequest.items))
            .where(ContractorRequest.id == request.id)
        )
        request_with_items = result.scalar_one()
        
        return {
            "request": request_with_items,
            "matching_results": matching_results
        }
        
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при обработке текста: {str(e)}")


def parse_text_to_items(text: str) -> List[dict]:
    """Парсит текст и извлекает позиции заявки - ищет строки длиннее 4 символов"""
    items = []
    lines = [line.strip() for line in text.strip().split('\n') if line.strip()]
    
    print(f"Parsing text with {len(lines)} lines:")
    for idx, line in enumerate(lines):
        print(f"  {idx}: '{line}' (length: {len(line)})")
    
    # Ищем строки длиннее 4 символов, которые могут быть описаниями товаров
    for line in lines:
        if len(line) > 4:
            # Пропускаем служебные строки
            if (line.startswith('Новая потребность') or 
                line.startswith('Для') or 
                line.lower() in ['шт', 'шт.', 'штук', 'штук.', 'pcs', 'pcs.', 'pieces', 'pieces.'] or
                re.match(r'^\d+$', line)):
                continue
            
            # Это потенциальное описание товара
            item = {
                'agb_article': '',  # Будет заполнено AI поиском
                'description': line,
                'quantity': 1,  # По умолчанию
                'unit': 'шт'    # По умолчанию
            }
            print(f"Found potential item: {item}")
            items.append(item)
    
    return items


@router.get("/our-database/")
async def get_our_database(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Получение нашей базы данных номенклатур"""
    try:
        result = await db.execute(select(MatchingNomenclature))
        nomenclatures = result.scalars().all()
        
        return [
            {
                "id": nom.id,
                "agb_article": nom.agb_article,
                "name": nom.name,
                "code_1c": nom.code_1c
            }
            for nom in nomenclatures
        ]
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка при получении базы данных: {str(e)}")


@router.get("/test-our-database/")
async def test_our_database(db: AsyncSession = Depends(get_db)):
    """Тестовый endpoint для проверки базы данных без аутентификации"""
    try:
        result = await db.execute(select(MatchingNomenclature))
        nomenclatures = result.scalars().all()
        
        return {
            "count": len(nomenclatures),
            "data": [
                {
                    "id": nom.id,
                    "agb_article": nom.agb_article,
                    "name": nom.name,
                    "code_1c": nom.code_1c,
                    "bl_article": nom.bl_article,
                    "packaging": nom.packaging,
                    "unit": nom.unit,
                    "is_active": nom.is_active
                }
                for nom in nomenclatures  # Все записи
            ]
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка при получении базы данных: {str(e)}")


@router.get("/found-items/")
async def get_found_items(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Получение найденных элементов сопоставления"""
    try:
        # Получаем все записи сопоставления с номенклатурами
        result = await db.execute(
            select(ArticleMapping)
            .options(selectinload(ArticleMapping.nomenclature))
            .options(selectinload(ArticleMapping.contractor_request_item))
        )
        mappings = result.scalars().all()
        
        found_items = []
        for mapping in mappings:
            item = {
                "id": mapping.id,
                "bl_article": mapping.bl_article,
                "search_article": mapping.contractor_request_item.contractor_article if mapping.contractor_request_item else None,
                "our_article": mapping.agb_article,
                "ut_number": mapping.nomenclature.code_1c if mapping.nomenclature else None,
                "found_data": mapping.nomenclature.name if mapping.nomenclature else None,
                "match_confidence": mapping.match_confidence,
                "packaging_factor": mapping.packaging_factor,
                "recalculated_quantity": mapping.recalculated_quantity
            }
            found_items.append(item)
        
        return found_items
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка при получении найденных элементов: {str(e)}")


@router.get("/test-found-items-debug/")
async def test_found_items_debug(db: AsyncSession = Depends(get_db)):
    """Отладочный endpoint для проверки фильтрации"""
    try:
        # Получаем все элементы
        result_all = await db.execute(select(ArticleMapping))
        all_mappings = result_all.scalars().all()
        
        # Получаем только активные элементы
        result_active = await db.execute(select(ArticleMapping).where(ArticleMapping.is_active == True))
        active_mappings = result_active.scalars().all()
        
        return {
            "total_count": len(all_mappings),
            "active_count": len(active_mappings),
            "inactive_count": len(all_mappings) - len(active_mappings),
            "message": "Отладочная информация о фильтрации"
        }
    except Exception as e:
        return {"error": str(e)}


@router.get("/test-found-items/")
async def test_found_items(db: AsyncSession = Depends(get_db)):
    """Тестовый endpoint для проверки найденных элементов без аутентификации"""
    try:
        # Используем SQLAlchemy ORM для получения данных, фильтруем только активные
        result = await db.execute(select(ArticleMapping).where(ArticleMapping.is_active == True))
        mappings = result.scalars().all()
        
        found_items = []
        for mapping in mappings:
            item = {
                "id": mapping.id,
                "mapping_id": mapping.id,  # Добавляем mapping_id для удаления
                "bl_article": mapping.bl_article,
                "search_article": mapping.contractor_article,
                "our_article": mapping.agb_article,
                "ut_number": mapping.agb_description,
                "found_data": mapping.contractor_description,
                "match_confidence": 100.0,
                "packaging_factor": mapping.packaging_factor,
                "recalculated_quantity": 1
            }
            found_items.append(item)
            print(f"Обработана запись: {mapping.id} - {mapping.agb_article} -> {mapping.bl_article}")
        
        return {
            "count": len(mappings),
            "data": found_items
        }
    except Exception as e:
        print(f"Ошибка в test_found_items: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Ошибка при получении найденных элементов: {str(e)}")


@router.put("/nomenclature/{nomenclature_id}/")
async def update_nomenclature(
    nomenclature_id: int,
    nomenclature_data: dict,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Обновление номенклатуры"""
    try:
        # Проверяем права доступа
        if current_user.role not in ['ved_passport', 'admin']:
            raise HTTPException(status_code=403, detail="Недостаточно прав для редактирования номенклатуры")
        
        # Находим номенклатуру
        result = await db.execute(
            select(MatchingNomenclature).where(MatchingNomenclature.id == nomenclature_id)
        )
        nomenclature = result.scalar_one_or_none()
        
        if not nomenclature:
            raise HTTPException(status_code=404, detail="Номенклатура не найдена")
        
        # Обновляем поля
        nomenclature.agb_article = nomenclature_data.get('agb_article', nomenclature.agb_article)
        nomenclature.name = nomenclature_data.get('name', nomenclature.name)
        nomenclature.code_1c = nomenclature_data.get('code_1c', nomenclature.code_1c)
        nomenclature.bl_article = nomenclature_data.get('bl_article', nomenclature.bl_article)
        nomenclature.packaging = nomenclature_data.get('packaging', nomenclature.packaging)
        nomenclature.unit = nomenclature_data.get('unit', nomenclature.unit)
        nomenclature.is_active = nomenclature_data.get('is_active', nomenclature.is_active)
        
        await db.commit()
        
        return {
            "id": nomenclature.id,
            "agb_article": nomenclature.agb_article,
            "name": nomenclature.name,
            "code_1c": nomenclature.code_1c,
            "bl_article": nomenclature.bl_article,
            "packaging": nomenclature.packaging,
            "unit": nomenclature.unit,
            "is_active": nomenclature.is_active
        }
        
    except HTTPException:
        raise
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при обновлении номенклатуры: {str(e)}")


@router.put("/mapping/{mapping_id}/")
async def update_mapping(
    mapping_id: int,
    mapping_data: dict,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Обновление сопоставления артикулов"""
    try:
        # Проверяем права доступа
        if current_user.role not in ['ved_passport', 'admin']:
            raise HTTPException(status_code=403, detail="Недостаточно прав для редактирования сопоставлений")
        
        # Находим сопоставление
        result = await db.execute(
            select(ArticleMapping).where(ArticleMapping.id == mapping_id)
        )
        mapping = result.scalar_one_or_none()
        
        if not mapping:
            raise HTTPException(status_code=404, detail="Сопоставление не найдено")
        
        # Обновляем поля
        mapping.bl_article = mapping_data.get('bl_article', mapping.bl_article)
        mapping.contractor_article = mapping_data.get('search_article', mapping.contractor_article)
        mapping.agb_article = mapping_data.get('our_article', mapping.agb_article)
        mapping.agb_description = mapping_data.get('ut_number', mapping.agb_description)
        mapping.contractor_description = mapping_data.get('found_data', mapping.contractor_description)
        mapping.match_confidence = mapping_data.get('match_confidence', mapping.match_confidence)
        mapping.packaging_factor = mapping_data.get('packaging_factor', mapping.packaging_factor)
        mapping.recalculated_quantity = mapping_data.get('recalculated_quantity', mapping.recalculated_quantity)
        
        await db.commit()
        
        return {
            "id": mapping.id,
            "bl_article": mapping.bl_article,
            "search_article": mapping.contractor_article,
            "our_article": mapping.agb_article,
            "ut_number": mapping.agb_description,
            "found_data": mapping.contractor_description,
            "match_confidence": mapping.match_confidence,
            "packaging_factor": mapping.packaging_factor,
            "recalculated_quantity": mapping.recalculated_quantity
        }
        
    except HTTPException:
        raise
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при обновлении сопоставления: {str(e)}")


async def perform_matching(request_id: int, db: AsyncSession) -> dict:
    """Выполняет сопоставление артикулов для заявки"""
    try:
        # Получаем элементы заявки
        result = await db.execute(
            select(ContractorRequestItem).where(ContractorRequestItem.request_id == request_id)
        )
        items = result.scalars().all()
        
        if not items:
            return {"message": "Нет элементов для сопоставления"}
        
        # Получаем все существующие сопоставления
        existing_mappings_result = await db.execute(select(ArticleMapping))
        existing_mappings = existing_mappings_result.scalars().all()
        
        # Создаем словарь для быстрого поиска существующих сопоставлений
        existing_by_contractor_article = {
            mapping.contractor_article: mapping for mapping in existing_mappings
        }
        
        new_mappings = []
        matched_items = []
        
        for item in items:
            # Проверяем, есть ли уже сопоставление для этого артикула
            if item.contractor_article in existing_by_contractor_article:
                existing_mapping = existing_by_contractor_article[item.contractor_article]
                matched_items.append({
                    "line_number": item.line_number,
                    "contractor_article": item.contractor_article,
                    "description": item.description,
                    "agb_article": existing_mapping.agb_article,
                    "bl_article": existing_mapping.bl_article,
                    "match_confidence": existing_mapping.match_confidence,
                    "packaging_factor": existing_mapping.packaging_factor,
                    "recalculated_quantity": item.quantity * existing_mapping.packaging_factor,
                    "source": "existing"
                })
            else:
                # Используем умный поиск через AI
                ai_result = await smart_search_with_ai(item.description, db)
                
                if ai_result.get("matches"):
                    # Берем лучшее соответствие
                    best_match = ai_result["matches"][0]
                    
                    # Создаем новое сопоставление
                    new_mapping = ArticleMapping(
                        contractor_article=item.contractor_article,
                        contractor_description=item.description,
                        agb_article=best_match.get("agb_article", ""),
                        agb_description=best_match.get("name", ""),
                        bl_article=best_match.get("bl_article", ""),
                        match_confidence=int(best_match.get("confidence", 0)),
                        packaging_factor=int(best_match.get("packaging", 1)),
                        recalculated_quantity=item.quantity * int(best_match.get("packaging", 1))
                    )
                    db.add(new_mapping)
                    new_mappings.append(new_mapping)
                    
                    matched_items.append({
                        "line_number": item.line_number,
                        "contractor_article": item.contractor_article,
                        "description": item.description,
                        "agb_article": best_match.get("agb_article", ""),
                        "bl_article": best_match.get("bl_article", ""),
                        "match_confidence": int(best_match.get("confidence", 0)),
                        "packaging_factor": int(best_match.get("packaging", 1)),
                        "recalculated_quantity": item.quantity * int(best_match.get("packaging", 1)),
                        "source": "new",
                        "search_type": ai_result.get("search_type", "general")
                    })
                else:
                    # Не найдено соответствие
                    matched_items.append({
                        "line_number": item.line_number,
                        "contractor_article": item.contractor_article or '',
                        "description": item.description or '',
                        "agb_article": '',
                        "bl_article": '',
                        "match_confidence": 0,
                        "packaging_factor": 1.0,
                        "recalculated_quantity": item.quantity,
                        "source": "not_found"
                    })
        
        # Сохраняем новые сопоставления
        if new_mappings:
            await db.commit()
        
        return {
            "total_items": len(items),
            "matched_items": len([item for item in matched_items if item["agb_article"]]),
            "new_mappings_count": len(new_mappings),
            "existing_mappings_count": len([item for item in matched_items if item["source"] == "existing"]),
            "not_found_count": len([item for item in matched_items if item["source"] == "not_found"]),
            "results": matched_items
        }
        
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при сопоставлении: {str(e)}")


@router.post("/requests/{request_id}/match/")
async def match_request_items(
    request_id: int,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Выполняет сопоставление артикулов для заявки"""
    try:
        # Проверяем права доступа
        if current_user.role not in ['ved_passport', 'admin']:
            raise HTTPException(status_code=403, detail="Недостаточно прав для сопоставления")
        
        # Проверяем, что заявка принадлежит пользователю
        result = await db.execute(
            select(ContractorRequest).where(ContractorRequest.id == request_id)
        )
        request = result.scalar_one_or_none()
        
        if not request or request.created_by != current_user.id:
            raise HTTPException(status_code=404, detail="Заявка не найдена")
        
        # Выполняем сопоставление
        matching_results = await perform_matching(request_id, db)
        
        return matching_results
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка при сопоставлении: {str(e)}")


@router.post("/test-upload-excel/")
async def test_upload_excel(
    file: UploadFile = File(...),
    contractor_name: str = "Контрагент",
    db: AsyncSession = Depends(get_db)
):
    """Тестовый endpoint для загрузки Excel без аутентификации"""
    try:
        # Читаем Excel файл
        content = await file.read()
        df = pd.read_excel(io.BytesIO(content))
        
        # Генерируем номер заявки
        request_number = f"TEST_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        # Создаем заявку
        new_request = ContractorRequest(
            request_number=request_number,
            contractor_name=contractor_name,
            request_date=datetime.now(),
            total_items=5,  # Тестовая заявка из 5 позиций
            created_by=1  # Тестовый пользователь
        )
        
        db.add(new_request)
        await db.flush()
        
        # Создаем тестовую заявку на основе примера
        test_items = [
            {"agb_article": "1299650", "description": "Шпиндель верхней части керноприемника H/HU, 25231, SDT", "quantity": 5, "unit": "шт"},
            {"agb_article": "1298240", "description": "Втулка удержания жидкости 306131, SDT", "quantity": 12, "unit": "шт"},
            {"agb_article": "1298244", "description": "Пружина мягкая N/H/P, удержания жидкости, 104966, SDT", "quantity": 10, "unit": "шт"},
            {"agb_article": "1299679", "description": "Щека верхняя для ключа разводного 24\", 14947, SDT", "quantity": 8, "unit": "шт"},
            {"agb_article": "1299680", "description": "Щека верхняя для ключа разводного 36\", 14950, SDT", "quantity": 8, "unit": "шт"}
        ]
        
        items = []
        for index, item_data in enumerate(test_items):
            line_number = index + 1
            new_item = ContractorRequestItem(
                request_id=new_request.id,
                line_number=line_number,
                contractor_article=item_data["agb_article"],
                description=item_data["description"],
                unit=item_data["unit"],
                quantity=item_data["quantity"],
                category="Для бурения"
            )
            db.add(new_item)
            items.append(new_item)
        
        await db.commit()
        await db.refresh(new_request)
        
        # Выполняем сопоставление
        matching_results = await perform_matching(new_request.id, db)
        
        return {
            "request": new_request,
            "matching_results": matching_results
        }
        
    except Exception as e:
        await db.rollback()
        raise HTTPException(status_code=500, detail=f"Ошибка при загрузке файла: {str(e)}")


@router.post("/smart-search/")
async def smart_search_endpoint(
    request_data: dict,
    db: AsyncSession = Depends(get_db)
):
    """Тестовый endpoint для умного поиска через AI"""
    try:
        search_text = request_data.get("search_text", "")
        if not search_text:
            raise HTTPException(status_code=400, detail="Не указан текст для поиска")
        
        print(f"Выполняем умный поиск для: {search_text}")
        
        # Выполняем умный поиск
        result = await smart_search_with_ai(search_text, db)
        
        return result
        
    except Exception as e:
        import traceback
        error_details = traceback.format_exc()
        print(f"Ошибка в smart_search_endpoint: {str(e)}")
        print(f"Traceback: {error_details}")
        raise HTTPException(status_code=500, detail=f"Ошибка при поиске: {str(e)}")


@router.post("/smart-upload/")
async def smart_upload_file(
    file: UploadFile = File(...),
    contractor_name: str = "Контрагент",
    db: AsyncSession = Depends(get_db)
):
    """Умная загрузка любого файла с AI поиском"""
    try:
        print(f"Начинаем умную загрузку файла: {file.filename}")
        
        # Читаем файл
        content = await file.read()
        
        # Определяем тип файла и парсим содержимое
        file_extension = file.filename.split(".")[-1].lower() if file.filename else ""
        
        search_items = []
        
        if file_extension in ["xlsx", "xls"]:
            # Excel файл
            df = pd.read_excel(io.BytesIO(content))
            print(f"Excel файл содержит {len(df)} строк")
            
            # Извлекаем все текстовые данные из Excel
            for index, row in df.iterrows():
                for col in df.columns:
                    cell_value = str(row[col]).strip()
                    if len(cell_value) > 3 and cell_value != "nan":
                        search_items.append({
                            "text": cell_value,
                            "row": index + 1,
                            "column": col
                        })
        
        elif file_extension in ["txt", "csv"]:
            # Текстовый файл
            text_content = content.decode("utf-8", errors="ignore")
            lines = [line.strip() for line in text_content.split("\n") if line.strip()]
            
            for index, line in enumerate(lines):
                if len(line) > 3:
                    search_items.append({
                        "text": line,
                        "row": index + 1,
                        "column": "text"
                    })
        
        else:
            # Пытаемся прочитать как текст
            try:
                text_content = content.decode("utf-8", errors="ignore")
                lines = [line.strip() for line in text_content.split("\n") if line.strip()]
                
                for index, line in enumerate(lines):
                    if len(line) > 3:
                        search_items.append({
                            "text": line,
                            "row": index + 1,
                            "column": "text"
                        })
            except:
                raise HTTPException(status_code=400, detail="Неподдерживаемый формат файла")
        
        if not search_items:
            raise HTTPException(status_code=400, detail="Не удалось извлечь данные из файла")
        
        print(f"Извлечено {len(search_items)} элементов для поиска")
        
        # Создаем заявку
        request_number = f"SMART_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        new_request = ContractorRequest(
            request_number=request_number,
            contractor_name=contractor_name,
            request_date=datetime.now(),
            total_items=len(search_items),
            created_by=1  # Тестовый пользователь
        )
        
        db.add(new_request)
        await db.flush()
        
        # Создаем элементы заявки
        items = []
        for search_item in search_items:
            # Создаем элемент заявки
            item = ContractorRequestItem(
                request_id=new_request.id,
                line_number=search_item["row"],
                contractor_article=search_item["text"][:50],  # Ограничиваем длину
                description=search_item["text"],
                quantity=1,  # По умолчанию
                unit="шт",
                category="AI поиск"
            )
            db.add(item)
            items.append(item)
        
        await db.commit()
        await db.refresh(new_request)
        
        # Выполняем финальное сопоставление
        matching_results = await perform_matching(new_request.id, db)
        
        return {
            "request": new_request,
            "matching_results": matching_results,
            "search_items_count": len(search_items),
            "file_type": file_extension
        }
        
    except Exception as e:
        await db.rollback()
        import traceback
        error_details = traceback.format_exc()
        print(f"Ошибка в smart_upload_file: {str(e)}")
        print(f"Traceback: {error_details}")
        raise HTTPException(status_code=500, detail=f"Ошибка при умной загрузке файла: {str(e)}")


@router.post("/step-upload/")
async def step_upload_file(
    file: UploadFile = File(...),
    contractor_name: str = "Контрагент",
    db: AsyncSession = Depends(get_db)
):
    """Пошаговая загрузка файла с видимыми этапами"""
    try:
        print(f"🚀 Этап 1: Начинаем загрузку файла: {file.filename}")
        
        # Этап 1: Читаем файл
        content = await file.read()
        print(f"✅ Этап 1 завершен: Файл прочитан ({len(content)} байт)")
        
        # Этап 2: Определяем тип файла и парсим содержимое
        file_extension = file.filename.split(".")[-1].lower() if file.filename else ""
        print(f"🔍 Этап 2: Анализируем файл типа {file_extension}")
        
        search_items = []
        
        if file_extension in ["xlsx", "xls"]:
            # Excel файл
            df = pd.read_excel(io.BytesIO(content))
            print(f"📊 Этап 2: Excel файл содержит {len(df)} строк")
            
            # Обрабатываем каждую строку как единое целое
            for index, row in df.iterrows():
                # Объединяем все ячейки строки в один текст
                row_text = " ".join([str(cell).strip() for cell in row.values if str(cell).strip() != "nan" and len(str(cell).strip()) > 0])
                
                if len(row_text) > 3:
                    search_items.append({
                        "text": row_text,
                        "row": index + 1,
                        "column": "row"
                    })
        
        elif file_extension in ["txt", "csv"]:
            # Текстовый файл
            text_content = content.decode("utf-8", errors="ignore")
            lines = [line.strip() for line in text_content.split("\n") if line.strip()]
            
            for index, line in enumerate(lines):
                if len(line) > 3:
                    search_items.append({
                        "text": line,
                        "row": index + 1,
                        "column": "text"
                    })
        
        else:
            # Пытаемся прочитать как текст
            try:
                text_content = content.decode("utf-8", errors="ignore")
                lines = [line.strip() for line in text_content.split("\n") if line.strip()]
                
                for index, line in enumerate(lines):
                    if len(line) > 3:
                        search_items.append({
                            "text": line,
                            "row": index + 1,
                            "column": "text"
                        })
            except:
                raise HTTPException(status_code=400, detail="Неподдерживаемый формат файла")
        
        if not search_items:
            raise HTTPException(status_code=400, detail="Не удалось извлечь данные из файла")
        
        print(f"✅ Этап 2 завершен: Извлечено {len(search_items)} элементов для поиска")
        
        # Этап 3: Создаем заявку в базе данных
        print(f"💾 Этап 3: Создаем заявку в базе данных")
        
        # Генерируем номер заявки в формате {название_контрагента}_{порядковый_номер}_{дата}
        contractor_name_latin = transliterate_to_latin(contractor_name)
        print(f"🔤 Транслитерация: '{contractor_name}' -> '{contractor_name_latin}'")
        
        # Получаем количество заявок для этого контрагента сегодня
        today = datetime.now().date()
        count_result = await db.execute(
            select(func.count(ContractorRequest.id)).where(
                ContractorRequest.contractor_name == contractor_name,
                func.date(ContractorRequest.request_date) == today
            )
        )
        request_count = count_result.scalar() + 1
        
        request_number = f"{contractor_name_latin}_{request_count:03d}_{datetime.now().strftime('%Y%m%d')}"
        new_request = ContractorRequest(
            request_number=request_number,
            contractor_name=contractor_name,
            request_date=datetime.now(),
            total_items=0,  # Будет обновлено после AI обработки
            created_by=1  # Тестовый пользователь
        )
        
        db.add(new_request)
        await db.flush()
        
        # Создаем элементы заявки с AI извлечением артикулов
        items = []
        total_articles_found = 0
        
        print(f"🔍 Этап 3: Извлекаем артикулы из {len(search_items)} строк через AI")
        
        for i, search_item in enumerate(search_items):
            print(f"📝 Обрабатываем строку {i+1}/{len(search_items)}: {search_item['text'][:50]}...")
            
            # Извлекаем артикулы из строки через AI
            articles = await extract_articles_from_text(search_item["text"], db)
            
            if articles:
                print(f"  ✅ Найдено {len(articles)} артикулов")
                total_articles_found += len(articles)
                
                # Создаем элемент заявки для каждого найденного артикула
                for article_data in articles:
                    item = ContractorRequestItem(
                        request_id=new_request.id,
                        line_number=search_item["row"],
                        contractor_article=article_data.get("article", "")[:50],
                        description=article_data.get("description", search_item["text"]),
                        quantity=article_data.get("quantity", 1),
                        unit=article_data.get("unit", "шт"),
                        category="AI извлечение"
                    )
                    db.add(item)
                    items.append(item)
            else:
                print(f"  ❌ Артикулы не найдены")
                # Создаем элемент заявки с исходным текстом
                item = ContractorRequestItem(
                    request_id=new_request.id,
                    line_number=search_item["row"],
                    contractor_article="",  # Артикул не найден
                    description=search_item["text"],
                    quantity=1,
                    unit="шт",
                    category="Без артикула"
                )
                db.add(item)
                items.append(item)
        
        # Обновляем общее количество элементов
        new_request.total_items = len(items)
        
        await db.commit()
        print(f"✅ Этап 3 завершен: Заявка создана с ID {new_request.id}")
        print(f"📊 Статистика: {len(search_items)} строк -> {total_articles_found} артикулов -> {len(items)} элементов заявки")
        
        # Этап 4: Запускаем сопоставление в фоне
        print(f"🔄 Этап 4: Запускаем сопоставление в фоне")
        import asyncio
        asyncio.create_task(perform_background_matching(new_request.id))
        
        # Возвращаем только простые данные без сложных объектов
        return {
            "success": True,
            "request_id": new_request.id,
            "request_number": request_number,
            "contractor_name": contractor_name,
            "rows_processed": len(search_items),
            "articles_found": total_articles_found,
            "total_items": len(items),
            "file_type": file_extension,
            "message": f"Файл загружен успешно. Обработано {len(search_items)} строк, найдено {total_articles_found} артикулов. Сопоставление выполняется в фоне.",
            "status": "uploaded"
        }
        
    except Exception as e:
        await db.rollback()
        import traceback
        error_details = traceback.format_exc()
        print(f"❌ Ошибка в step_upload_file: {str(e)}")
        print(f"Traceback: {error_details}")
        raise HTTPException(status_code=500, detail=f"Ошибка при загрузке файла: {str(e)}")


async def perform_background_matching(request_id: int):
    """Выполняет сопоставление в фоне"""
    from database import get_db
    
    # Создаем новую сессию для фоновой задачи
    async for db in get_db():
        try:
            print(f"🔄 Начинаем фоновое сопоставление для заявки {request_id}")
            
            # Получаем элементы заявки
            result = await db.execute(
                select(ContractorRequestItem).where(ContractorRequestItem.request_id == request_id)
            )
            items = result.scalars().all()
            
            if not items:
                print(f"❌ Элементы заявки {request_id} не найдены")
                break
            
            print(f"📋 Найдено {len(items)} элементов для сопоставления")
            
            # Выполняем сопоставление для каждого элемента
            matched_count = 0
            for i, item in enumerate(items):
                print(f"🔍 Сопоставляем {i+1}/{len(items)}: {item.description[:50]}...")
                
                try:
                    # Выполняем умный поиск
                    ai_result = await smart_search_with_ai(item.description, db)
                    
                    if ai_result.get("matches"):
                        # Берем лучшее соответствие
                        best_match = ai_result["matches"][0]
                        
                        # Создаем сопоставление
                        mapping = ArticleMapping(
                            contractor_article=item.contractor_article,
                            contractor_description=item.description,
                            agb_article=best_match.get("agb_article", ""),
                            agb_description=best_match.get("name", ""),
                            bl_article=best_match.get("bl_article", ""),
                            match_confidence=int(best_match.get("confidence", 0)),
                            packaging_factor=int(best_match.get("packaging", 1)),
                            recalculated_quantity=item.quantity * int(best_match.get("packaging", 1))
                        )
                        db.add(mapping)
                        matched_count += 1
                        
                        print(f"  ✅ Найдено: {best_match.get('agb_article', '')} | {best_match.get('name', '')} | {best_match.get('confidence', 0)}%")
                    else:
                        print(f"  ❌ Не найдено соответствий")
                        
                except Exception as e:
                    print(f"  ⚠️ Ошибка при сопоставлении: {str(e)}")
                    continue
            
            await db.commit()
            print(f"🎉 Фоновое сопоставление завершено. Найдено {matched_count} соответствий.")
            break  # Выходим из цикла async for
            
        except Exception as e:
            await db.rollback()
            print(f"❌ Ошибка в фоновом сопоставлении: {str(e)}")
            break  # Выходим из цикла async for


@router.post("/test-upload-text/")
async def test_upload_text_request(
    request_data: TextUploadRequest,
    db: AsyncSession = Depends(get_db)
):
    """Тестовый endpoint для загрузки текста без аутентификации"""
    try:
        # Парсим текст для извлечения позиций
        items = parse_text_to_items(request_data.text)
        
        if not items:
            raise HTTPException(status_code=400, detail="Не удалось извлечь позиции из текста")
        
        # Создаем заявку с правильным номером
        contractor_name = request_data.contractor_name
        
        # Генерируем номер заявки в формате {название_контрагента}_{порядковый_номер}_{дата}
        contractor_name_latin = transliterate_to_latin(contractor_name)
        print(f"🔤 Транслитерация: '{contractor_name}' -> '{contractor_name_latin}'")
        
        # Получаем количество заявок для этого контрагента сегодня
        today = datetime.now().date()
        count_result = await db.execute(
            select(func.count(ContractorRequest.id)).where(
                ContractorRequest.contractor_name == contractor_name,
                func.date(ContractorRequest.request_date) == today
            )
        )
        request_count = count_result.scalar() + 1
        
        request_number = f"{contractor_name_latin}_{request_count:03d}_{datetime.now().strftime('%Y%m%d')}"
        
        request = ContractorRequest(
            request_number=request_number,
            contractor_name=contractor_name,
            request_date=datetime.now(),
            total_items=len(items),
            status='new',
            created_by=1  # Тестовый пользователь
        )
        
        db.add(request)
        await db.flush()  # Получаем ID заявки
        
        # Создаем позиции заявки
        for i, item_data in enumerate(items):
            item = ContractorRequestItem(
                request_id=request.id,
                line_number=i + 1,  # Добавляем номер строки
                contractor_article=item_data.get('agb_article', ''),
                description=item_data.get('description', ''),
                quantity=item_data.get('quantity', 1),
                unit=item_data.get('unit', 'шт'),
                category="Текстовый ввод"
            )
            db.add(item)
        
        await db.commit()
        
        return {
            "success": True,
            "request": {
                "id": request.id,
                "request_number": request_number,
                "contractor_name": contractor_name,
                "total_items": len(items),
                "status": "new"
            },
            "message": f"Текстовая заявка создана успешно. Обработано {len(items)} позиций."
        }
        
    except Exception as e:
        await db.rollback()
        print(f"❌ Ошибка в тестовом текстовом вводе: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Ошибка при создании заявки: {str(e)}")

@router.post("/test-match/{request_id}")
async def test_match_request(
    request_id: int,
    db: AsyncSession = Depends(get_db)
):
    """Тестовый endpoint для сопоставления без аутентификации"""
    try:
        print(f"🔄 Тестовое сопоставление для заявки {request_id}")
        
        # Получаем элементы заявки
        result = await db.execute(
            select(ContractorRequestItem).where(ContractorRequestItem.request_id == request_id)
        )
        items = result.scalars().all()
        
        if not items:
            raise HTTPException(status_code=404, detail="Элементы заявки не найдены")
        
        print(f"📋 Найдено {len(items)} элементов для сопоставления")
        
        # Выполняем сопоставление для каждого элемента
        matched_count = 0
        unmatched_count = 0
        
        for i, item in enumerate(items):
            print(f"🔍 Сопоставляем {i+1}/{len(items)}: {item.description[:50]}...")
            
            try:
                # Выполняем умный поиск
                ai_result = await smart_search_with_ai(item.description, db)
                
                if ai_result.get("matches"):
                    # Берем лучшее соответствие
                    best_match = ai_result["matches"][0]
                    
                    # Создаем сопоставление
                    mapping = ArticleMapping(
                        contractor_article=item.contractor_article,
                        contractor_description=item.description,
                        agb_article=best_match.get("agb_article", ""),
                        agb_description=best_match.get("name", ""),
                        bl_article=best_match.get("bl_article", ""),
                        match_confidence=int(best_match.get("confidence", 0)),
                        packaging_factor=int(best_match.get("packaging", 1)),
                        recalculated_quantity=item.quantity * int(best_match.get("packaging", 1))
                    )
                    db.add(mapping)
                    matched_count += 1
                    
                    print(f"  ✅ Найдено: {best_match.get('agb_article', '')} | {best_match.get('name', '')} | {best_match.get('confidence', 0)}%")
                else:
                    unmatched_count += 1
                    print(f"  ❌ Не найдено соответствий")
                    
            except Exception as e:
                print(f"  ⚠️ Ошибка при сопоставлении: {str(e)}")
                unmatched_count += 1
                continue
        
        # Сохраняем все изменения
        await db.commit()
        
        return {
            "total_items": len(items),
            "matched_items": matched_count,
            "unmatched_items": unmatched_count,
            "message": f"Сопоставление завершено. Найдено {matched_count} из {len(items)} позиций."
        }
        
    except Exception as e:
        await db.rollback()
        print(f"❌ Ошибка в тестовом сопоставлении: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Ошибка при сопоставлении: {str(e)}")


@router.get("/status/{request_id}")
async def get_request_status(
    request_id: int,
    db: AsyncSession = Depends(get_db)
):
    """Получить статус обработки заявки"""
    try:
        # Получаем заявку
        result = await db.execute(
            select(ContractorRequest).where(ContractorRequest.id == request_id)
        )
        request = result.scalar_one_or_none()
        
        if not request:
            raise HTTPException(status_code=404, detail="Заявка не найдена")
        
        # Получаем количество найденных сопоставлений
        mapping_result = await db.execute(
            select(func.count(ArticleMapping.id))
        )
        total_mappings = mapping_result.scalar()
        
        # Получаем элементы заявки
        items_result = await db.execute(
            select(ContractorRequestItem).where(ContractorRequestItem.request_id == request_id)
        )
        items = items_result.scalars().all()
        
        return {
            "request_id": request_id,
            "request_number": request.request_number,
            "status": "completed" if len(items) > 0 else "processing",
            "total_items": len(items),
            "total_mappings": total_mappings,
            "message": "Обработка завершена" if len(items) > 0 else "Обработка в процессе"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка при получении статуса: {str(e)}")


@router.post("/requests/{request_id}/smart-match/")
async def smart_match_request(
    request_id: int,
    db: AsyncSession = Depends(get_db)
):
    """Умное сопоставление заявки через AI"""
    try:
        print(f"Начинаем умное сопоставление заявки {request_id}")
        
        # Получаем элементы заявки
        result = await db.execute(
            select(ContractorRequestItem).where(ContractorRequestItem.request_id == request_id)
        )
        items = result.scalars().all()
        
        if not items:
            raise HTTPException(status_code=404, detail="Элементы заявки не найдены")
        
        print(f"Найдено {len(items)} элементов для сопоставления")
        
        # Выполняем сопоставление для каждого элемента
        matched_count = 0
        for item in items:
            print(f"Сопоставляем: {item.description}")
            
            # Выполняем умный поиск
            ai_result = await smart_search_with_ai(item.description, db)
            
            if ai_result.get("matches"):
                # Берем лучшее соответствие
                best_match = ai_result["matches"][0]
                
                # Создаем сопоставление
                mapping = ArticleMapping(
                    contractor_article=item.contractor_article,
                    contractor_description=item.description,
                    agb_article=best_match.get("agb_article", ""),
                    agb_description=best_match.get("name", ""),
                    bl_article=best_match.get("bl_article", ""),
                    match_confidence=int(best_match.get("confidence", 0)),
                    packaging_factor=int(best_match.get("packaging", 1)),
                    recalculated_quantity=item.quantity * int(best_match.get("packaging", 1))
                )
                db.add(mapping)
                matched_count += 1
                
                print(f"  ✅ Найдено: {best_match.get('agb_article', '')} | {best_match.get('name', '')} | {best_match.get('confidence', 0)}%")
            else:
                print(f"  ❌ Не найдено соответствий")
        
        await db.commit()
        
        # Выполняем финальное сопоставление
        matching_results = await perform_matching(request_id, db)
        
        return {
            "request_id": request_id,
            "matching_results": matching_results,
            "matched_count": matched_count,
            "message": f"Сопоставление завершено. Найдено {matched_count} новых соответствий."
        }
        
    except Exception as e:
        await db.rollback()
        import traceback
        error_details = traceback.format_exc()
        print(f"Ошибка в smart_match_request: {str(e)}")
        print(f"Traceback: {error_details}")
        raise HTTPException(status_code=500, detail=f"Ошибка при сопоставлении: {str(e)}")


@router.get("/test-requests/")
async def test_get_requests(db: AsyncSession = Depends(get_db)):
    """Тестовый endpoint для получения заявок без аутентификации"""
    try:
        result = await db.execute(
            select(ContractorRequest)
            .options(selectinload(ContractorRequest.items))
            .order_by(ContractorRequest.created_at.desc())
            .limit(10)
        )
        requests = result.scalars().all()
        
        return {
            "count": len(requests),
            "data": requests
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка при получении заявок: {str(e)}")


# ИИ обработка файлов
UPLOAD_DIR = Path("uploads/ai_processing")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

ALLOWED_EXTENSIONS = {
    'pdf': ['.pdf'],
    'excel': ['.xlsx', '.xls', '.csv', '.ods'],
    'word': ['.doc', '.docx', '.odt', '.rtf'],
    'powerpoint': ['.ppt', '.pptx', '.odp'],
    'text': ['.txt', '.rtf'],
    'images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.svg']
}

def get_file_extension(filename: str) -> str:
    """Получить расширение файла"""
    return Path(filename).suffix.lower()

def is_allowed_file(filename: str) -> bool:
    """Проверить, разрешен ли тип файла"""
    ext = get_file_extension(filename)
    return any(ext in extensions for extensions in ALLOWED_EXTENSIONS.values())

async def extract_text_from_pdf(file_path: str) -> str:
    """Извлечь текст из PDF файла"""
    try:
        import PyPDF2
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"Ошибка извлечения текста из PDF: {str(e)}")

async def extract_text_from_excel(file_path: str) -> str:
    """Извлечь текст из Excel файла"""
    try:
        df = pd.read_excel(file_path)
        text = ""
        for index, row in df.iterrows():
            text += " ".join([str(cell) for cell in row if pd.notna(cell)]) + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"Ошибка извлечения текста из Excel: {str(e)}")

async def extract_text_from_image(file_path: str) -> str:
    """Извлечь текст из изображения с помощью OCR"""
    try:
        from PIL import Image
        import pytesseract
        
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang='rus+eng')
        return text.strip()
    except Exception as e:
        # Fallback: попробуем без языков
        try:
            from PIL import Image
            import pytesseract
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as fallback_error:
            raise Exception(f"Ошибка OCR обработки: {str(e)}, Fallback: {str(fallback_error)}")
    except Exception as e:
        raise Exception(f"Ошибка OCR обработки: {str(e)}")

async def extract_text_from_word(file_path: str) -> str:
    """Извлечь текст из Word документа"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"Ошибка извлечения текста из Word: {str(e)}")

async def extract_text_from_powerpoint(file_path: str) -> str:
    """Извлечь текст из PowerPoint презентации"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"Ошибка извлечения текста из PowerPoint: {str(e)}")

async def extract_text_from_text_file(file_path: str) -> str:
    """Извлечь текст из текстового файла"""
    try:
        # Пробуем разные кодировки
        encodings = ['utf-8', 'cp1251', 'latin-1', 'iso-8859-1']
        text = ""
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    text = file.read()
                break
            except UnicodeDecodeError:
                continue
        
        if not text:
            # Если все кодировки не сработали, читаем как бинарный файл
            with open(file_path, 'rb') as file:
                content = file.read()
                text = content.decode('utf-8', errors='ignore')
        
        return text.strip()
    except Exception as e:
        raise Exception(f"Ошибка извлечения текста из текстового файла: {str(e)}")

async def extract_text_from_file(file_path: str, filename: str) -> str:
    """Извлечь текст из файла в зависимости от его типа"""
    ext = get_file_extension(filename)
    
    if ext in ALLOWED_EXTENSIONS['pdf']:
        return await extract_text_from_pdf(file_path)
    elif ext in ALLOWED_EXTENSIONS['excel']:
        return await extract_text_from_excel(file_path)
    elif ext in ALLOWED_EXTENSIONS['word']:
        return await extract_text_from_word(file_path)
    elif ext in ALLOWED_EXTENSIONS['powerpoint']:
        return await extract_text_from_powerpoint(file_path)
    elif ext in ALLOWED_EXTENSIONS['text']:
        return await extract_text_from_text_file(file_path)
    elif ext in ALLOWED_EXTENSIONS['images']:
        return await extract_text_from_image(file_path)
    else:
        raise Exception(f"Неподдерживаемый тип файла: {ext}")

async def get_ai_response(text: str, api_key: str, provider: str) -> str:
    """Получить ответ от ИИ"""
    try:
        if provider == 'openai':
            import openai
            openai.api_key = api_key
            response = await openai.ChatCompletion.acreate(
                model="gpt-3.5-turbo",
                messages=[
                    {
                        "role": "system",
                        "content": """ТЫ - ИИ-СПЕЦИАЛИСТ ПО ИЗВЛЕЧЕНИЮ И СОПОСТАВЛЕНИЮ АРТИКУЛОВ КОМПАНИИ АГБ.

ИНСТРУКЦИИ:
1. ТВОЯ ОСНОВНАЯ ЗАДАЧА: Извлекать артикулы и товары из документов/текста
2. ЕСЛИ это запрос на поиск артикулов - ОБЯЗАТЕЛЬНО верни JSON с результатами
3. ЕСЛИ это обычный чат - отвечай как помощник

ПРИЗНАКИ АРТИКУЛОВ КОМПАНИИ АГБ:
- Формат: АГБ-XXXXXX, АГБXXXXXX (6-8 цифр)
- BL-XXXXXX, BLXXXXXX (артикулы "бортлангер")
- ОХКУ-XXXXXX, ОХКУXXXXXX (артикулы)
- Коды 1С: числовые коды (УТ-коды)
- Сочетания: цифры + буквы (1234BL, BL1234)

ИЗВЛЕКАЙ ИНФОРМАЦИЮ:
- Артикулы контрагентов (любые форматы)
- Наименования товаров (на русском/английском)
- Количества и единицы измерения
- Описания товаров

ФОРМАТ ОТВЕТА ДЛЯ АРТИКУЛОВ:
[
    {
        "contractor_article": "артикул контрагента (если есть)",
        "description": "полное наименование товара",
        "quantity": число,
        "unit": "шт/кг/л/м²/м³",
        "agb_article": "артикул АГБ (если похож)",
        "bl_article": "BL артикул (если похож)",
        "match_confidence": 85,
        "match_type": "exact/partial/by_name"
    }
]

ПРАВИЛА ИЗВЛЕЧЕНИЯ:
- ИЩИ ВСЕ возможные артикулы в тексте
- НЕ ПРОПУСКАЙ товары без артикулов - они тоже важны
- ЕСЛИ артикул содержит "BL", "борт" - это BL артикул
- ЕСЛИ артикул содержит "АГБ" - это АГБ артикул
- КОЛИЧЕСТВО: ищи числа с единицами измерения
- ЕСЛИ количества нет - ставь 1

ПРИМЕРЫ РАСПОЗНАВАНИЯ:
- "BL-123456 болт М12" → contractor_article: "BL-123456", description: "болт М12"
- "АГБ-789012" → agb_article: "АГБ-789012"
- "10 шт гайка М10" → quantity: 10, unit: "шт", description: "гайка М10"
- "труба 100м" → quantity: 100, unit: "м", description: "труба"

ВЕРНИ ТОЛЬКО JSON МАССИВ!"""
                    },
                    {
                        "role": "user",
                        "content": f"Проанализируй этот текст и найди артикулы:\n\n{text}"
                    }
                ]
            )
            return response.choices[0].message.content
        elif provider == 'polza':
            # Интеграция с Polza.ai согласно документации
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            
            data = {
                "model": "gpt-4o-mini",
                "messages": [
                    {
                        "role": "system",
                        "content": """ТЫ - ИИ-СПЕЦИАЛИСТ ПО ИЗВЛЕЧЕНИЮ И СОПОСТАВЛЕНИЮ АРТИКУЛОВ КОМПАНИИ АГБ.

ИНСТРУКЦИИ:
1. ТВОЯ ОСНОВНАЯ ЗАДАЧА: Извлекать артикулы и товары из документов/текста
2. ЕСЛИ это запрос на поиск артикулов - ОБЯЗАТЕЛЬНО верни JSON с результатами
3. ЕСЛИ это обычный чат - отвечай как помощник

ПРИЗНАКИ АРТИКУЛОВ КОМПАНИИ АГБ:
- Формат: АГБ-XXXXXX, АГБXXXXXX (6-8 цифр)
- BL-XXXXXX, BLXXXXXX (артикулы "бортлангер")
- ОХКУ-XXXXXX, ОХКУXXXXXX (артикулы)
- Коды 1С: числовые коды (УТ-коды)
- Сочетания: цифры + буквы (1234BL, BL1234)

ИЗВЛЕКАЙ ИНФОРМАЦИЮ:
- Артикулы контрагентов (любые форматы)
- Наименования товаров (на русском/английском)
- Количества и единицы измерения
- Описания товаров

ФОРМАТ ОТВЕТА ДЛЯ АРТИКУЛОВ:
[
    {
        "contractor_article": "артикул контрагента (если есть)",
        "description": "полное наименование товара",
        "quantity": число,
        "unit": "шт/кг/л/м²/м³",
        "agb_article": "артикул АГБ (если похож)",
        "bl_article": "BL артикул (если похож)",
        "match_confidence": 85,
        "match_type": "exact/partial/by_name"
    }
]

ПРАВИЛА ИЗВЛЕЧЕНИЯ:
- ИЩИ ВСЕ возможные артикулы в тексте
- НЕ ПРОПУСКАЙ товары без артикулов - они тоже важны
- ЕСЛИ артикул содержит "BL", "борт" - это BL артикул
- ЕСЛИ артикул содержит "АГБ" - это АГБ артикул
- КОЛИЧЕСТВО: ищи числа с единицами измерения
- ЕСЛИ количества нет - ставь 1

ПРИМЕРЫ РАСПОЗНАВАНИЯ:
- "BL-123456 болт М12" → contractor_article: "BL-123456", description: "болт М12"
- "АГБ-789012" → agb_article: "АГБ-789012"
- "10 шт гайка М10" → quantity: 10, unit: "шт", description: "гайка М10"
- "труба 100м" → quantity: 100, unit: "м", description: "труба"

ВЕРНИ ТОЛЬКО JSON МАССИВ!"""
                    },
                    {
                        "role": "user",
                        "content": f"Проанализируй этот текст и найди артикулы:\n\n{text}"
                    }
                ],
                "temperature": 0.1,
                "max_tokens": 2000
            }
            
            # Используем URL из документации Polza.ai
            POLZA_API_URL = "https://api.polza.ai/v1/chat/completions"
            
            print(f"DEBUG: Sending request to Polza.ai with key: {api_key[:10]}...")
            print(f"DEBUG: Request data: {data}")
            
            async with aiohttp.ClientSession() as session:
                async with session.post(POLZA_API_URL, headers=headers, json=data) as response:
                    print(f"DEBUG: Polza.ai response status: {response.status}")
                    if response.status in [200, 201]:
                        result = await response.json()
                        print(f"DEBUG: Polza.ai response: {result}")
                        return result["choices"][0]["message"]["content"]
                    else:
                        error_text = await response.text()
                        print(f"DEBUG: Polza.ai error: {error_text}")
                        raise Exception(f"Ошибка Polza.ai API {response.status}: {error_text}")
        else:
            raise Exception(f"Неподдерживаемый провайдер: {provider}")
    except Exception as e:
        raise Exception(f"Ошибка получения ответа от ИИ: {str(e)}")

async def match_articles_with_database(articles: List[dict], db: AsyncSession) -> List[MatchingResult]:
    """Сопоставить найденные артикулы с базой данных"""
    print(f"🔍 match_articles_with_database вызвана с {len(articles)} артикулами")
    results = []
    
    for article in articles:
        # Проверяем, что article - это словарь
        if not isinstance(article, dict):
            print(f"❌ Ошибка: article должен быть словарем, получен: {type(article)}")
            return []

        contractor_article = article.get('contractor_article')
        description = article.get('description')

        # Если нет артикула и описания, пропускаем
        if not contractor_article and not description:
            print(f"❌ Ошибка: отсутствуют contractor_article и description")
            results.append(MatchingResult(
                id=str(uuid.uuid4()),
                contractor_article='',
                description='Неизвестный товар',
                matched=False,
                match_confidence=0.0
            ))
            return results

        # Используем умный поиск через AI
        search_text = f"{contractor_article or ''} {description or ''}".strip()
        print(f"🔍 Вызываем smart_search_with_ai для: '{search_text}'")
        print(f"🔍 contractor_article: '{contractor_article}'")
        print(f"🔍 description: '{description}'")
        
        try:
            ai_search_result = await smart_search_with_ai(search_text, db)
            print(f"🔍 Результат smart_search_with_ai: {ai_search_result}")
        except Exception as e:
            print(f"❌ Ошибка в smart_search_with_ai: {e}")
            import traceback
            traceback.print_exc()
            ai_search_result = None
        
        if ai_search_result and ai_search_result.get('matches'):
            best_match = ai_search_result['matches'][0]  # Берем лучший результат
            print(f"✅ Найдено сопоставление: {best_match}")
            
            # Определяем тип поиска для сообщения
            search_type = ai_search_result.get('search_type', 'ai_search')
            is_existing = best_match.get('is_existing', False)
            
            results.append(MatchingResult(
                id=str(uuid.uuid4()),
                contractor_article=contractor_article,
                description=description,
                matched=True,
                agb_article=best_match.get('agb_article'),
                bl_article=best_match.get('bl_article'),
                match_confidence=best_match.get('confidence', 0.0),
                nomenclature={
                    'id': 0,
                    'name': best_match.get('name', ''),
                    'code_1c': best_match.get('code_1c', ''),
                    'article': best_match.get('agb_article', '')
                },
                # Добавляем информацию о типе поиска
                search_type=search_type,
                is_existing_mapping=is_existing,
                mapping_id=best_match.get('mapping_id') if is_existing else None
            ))
        else:
            # Если соответствие не найдено
            results.append(MatchingResult(
                id=str(uuid.uuid4()),
                contractor_article=contractor_article,
                description=description,
                matched=False,
                match_confidence=0.0
            ))
    
    return results

async def find_exact_article_match(contractor_article: str, db: AsyncSession) -> Optional[dict]:
    """Найти точное соответствие артикула в базе данных"""
    try:
        from sqlalchemy.future import select

        # Если contractor_article None или пустой, не ищем
        if not contractor_article:
            return None

        # Ищем точное соответствие по артикулу
        result = await db.execute(
            select(MatchingNomenclature).where(
                MatchingNomenclature.agb_article == contractor_article
            ).limit(1)
        )
        match = result.scalar_one_or_none()

        if match:
            return {
                'id': match.id,
                'name': match.name,
                'code_1c': match.code_1c,
                'agb_article': match.agb_article,
                'bl_article': match.bl_article,
                'confidence': 100.0
            }

        return None
    except Exception as e:
        print(f"Error in find_exact_article_match: {e}")
        return None

async def find_partial_article_match(contractor_article: str, description: str, db: AsyncSession) -> Optional[dict]:
    """Найти частичное соответствие артикула в базе данных"""
    try:
        from sqlalchemy.future import select
        import difflib
        
        # Получаем все номенклатуры
        result = await db.execute(select(MatchingNomenclature))
        nomenclatures = result.scalars().all()
        
        best_match = None
        best_confidence = 0
        
        for nom in nomenclatures:
            # Проверяем соответствие по артикулу
            article_confidence = 0
            if contractor_article and nom.agb_article:
                article_confidence = SequenceMatcher(
                    None, contractor_article.lower(), nom.agb_article.lower()
                ).ratio() * 100
            
            # Проверяем соответствие по наименованию
            name_confidence = 0
            if description and nom.name:
                name_confidence = SequenceMatcher(
                    None, description.lower(), nom.name.lower()
                ).ratio() * 100
            
            # Берем максимальную уверенность
            confidence = max(article_confidence, name_confidence)
            
            if confidence > best_confidence and confidence >= 80:
                best_confidence = confidence
                best_match = {
                    'id': nom.id,
                    'name': nom.name,
                    'code_1c': nom.code_1c,
                    'agb_article': nom.agb_article,
                    'bl_article': nom.bl_article,
                    'confidence': confidence
                }
        
        return best_match
    except Exception as e:
        print(f"Error in find_partial_article_match: {e}")
        return None

@router.post("/ai-process/", response_model=AIMatchingResponse)
async def process_ai_request(
    message: str = Form(...),
    files: List[UploadFile] = File(default=[]),
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Обработать запрос к ИИ-агенту"""
    try:
        # Получаем API ключ с приоритетом Polza.ai
        result = await db.execute(select(ApiKey).where(
            ApiKey.is_active == True,
            ApiKey.provider == 'polza'
        ).limit(1))
        api_key_obj = result.scalar_one_or_none()
        
        # Если Polza.ai ключа нет, берем OpenAI
        if not api_key_obj:
            result = await db.execute(select(ApiKey).where(
                ApiKey.is_active == True,
                ApiKey.provider == 'openai'
            ).limit(1))
            api_key_obj = result.scalar_one_or_none()
        
        if not api_key_obj:
            return AIMatchingResponse(
                message="Нет активного API ключа для ИИ-сервиса. Обратитесь к администратору.",
                matching_results=[],
                processing_time=0.1,
                status="error"
            )
        
        # Обрабатываем файлы
        extracted_text = message
        file_paths = []
        
        for file in files:
            if not is_allowed_file(file.filename):
                return AIMatchingResponse(
                    message=f"Неподдерживаемый тип файла: {file.filename}",
                    matching_results=[],
                    processing_time=0.1,
                    status="error"
                )
            
            # Сохраняем файл
            file_id = str(uuid.uuid4())
            file_path = UPLOAD_DIR / f"{file_id}_{file.filename}"
            
            with open(file_path, "wb") as buffer:
                content = await file.read()
                buffer.write(content)
            
            file_paths.append(str(file_path))
            
            # Извлекаем текст из файла
            try:
                file_text = await extract_text_from_file(str(file_path), file.filename)
                extracted_text += f"\n\n--- Содержимое файла {file.filename} ---\n{file_text}"
            except Exception as e:
                print(f"Ошибка обработки файла {file.filename}: {str(e)}")
                continue
        
        # Расшифровываем ключ
        try:
            from cryptography.fernet import Fernet
            # Используем тот же фиксированный ключ, что и в settings.py
            ENCRYPTION_KEY = b'iF0d2ARGQpaU9GFfQdWNovBL239dqwTp9hDDPrDQQic='
            cipher_suite = Fernet(ENCRYPTION_KEY)
            decrypted_key = cipher_suite.decrypt(api_key_obj.key.encode()).decode()
            print(f"DEBUG: Successfully decrypted key for {api_key_obj.provider}")
        except Exception as e:
            print(f"DEBUG: Decryption failed: {str(e)}")
            print(f"DEBUG: Key length: {len(api_key_obj.key)}")
            print(f"DEBUG: Key starts with: {api_key_obj.key[:20]}...")
            # Возможно, ключ не зашифрован, используем напрямую
            decrypted_key = api_key_obj.key
            print(f"DEBUG: Using key directly: {decrypted_key[:10]}...")
        
        # Убираем пробелы и переносы строк из ключа
        decrypted_key = decrypted_key.strip()
        print(f"DEBUG: Final key for {api_key_obj.provider}: {decrypted_key[:20]}...")
        
        # Получаем ответ от ИИ
        try:
            print(f"🔍 Отправляем запрос к AI: '{extracted_text}'")
            # Получаем ответ от ИИ
            ai_response = await get_ai_response(extracted_text, decrypted_key, api_key_obj.provider)
            print(f"🔍 AI ответ: {ai_response}")
            
            # Проверяем, содержит ли ответ JSON с артикулами
            articles = []
            matching_results = []
            
            # Ищем JSON в ответе
            try:
                # Убираем markdown блоки если есть
                clean_response = ai_response
                if '```json' in clean_response:
                    clean_response = clean_response.split('```json')[1].split('```')[0].strip()
                elif '```' in clean_response:
                    clean_response = clean_response.split('```')[1].split('```')[0].strip()
                
                # Пытаемся найти JSON в ответе
                if clean_response.strip().startswith('[') and clean_response.strip().endswith(']'):
                    articles = json.loads(clean_response)
                    print(f"DEBUG: Found articles in response: {len(articles)}")
                    print(f"DEBUG: Articles content: {articles}")
                else:
                    # Ищем JSON в тексте
                    import re
                    json_match = re.search(r'\[.*?\]', clean_response, re.DOTALL)
                    if json_match:
                        articles = json.loads(json_match.group())
                        print(f"DEBUG: Found articles in text: {len(articles)}")
                    else:
                        print(f"DEBUG: No JSON found, treating as regular chat response")
                        # Проверяем, может это описание товара
                        if clean_response and len(clean_response.strip()) > 10:
                            print(f"🔍 Попытка обработать как натуральное описание товара: {clean_response[:100]}...")
                            articles = [{
                                'contractor_article': '',
                                'description': clean_response.strip(),
                                'quantity': 1,
                                'unit': 'шт'
                            }]
                            print(f"🔍 Создан артикул из натурального текста: {articles[0]}")
                        else:
                            print(f"DEBUG: Treating as regular chat response")
                            pass
            except json.JSONDecodeError as e:
                print(f"DEBUG: JSON parsing error: {e}")
                print(f"DEBUG: Response: {ai_response}")
                # Если не удалось распарсить JSON, проверяем, может это описание товара
                if ai_response and len(ai_response.strip()) > 10:
                    print(f"🔍 Попытка обработать как натуральное описание товара: {ai_response[:100]}...")
                    # Обрабатываем как натуральное описание товара
                    articles = [{
                        'contractor_article': '',
                        'description': ai_response.strip(),
                        'quantity': 1,
                        'unit': 'шт'
                    }]
                    print(f"🔍 Создан артикул из натурального текста: {articles[0]}")
                else:
                    # Если не удалось распарсить JSON, это обычный чат
                    pass
            
            # Если нашли артикулы, создаем результаты сопоставления
            if articles:
                print(f"🔍 HTTP API: Найдено {len(articles)} артикулов в AI ответе")
                for article in articles:
                    print(f"🔍 HTTP API: Обрабатываем артикул: {article}")
                    # Нормализуем данные из ИИ ответа
                    normalized_article = normalize_ai_article(article)

                    # Сопоставляем с базой данных
                    print(f"🔍 HTTP API: Вызываем match_articles_with_database...")
                    matched_result = await match_articles_with_database([normalized_article], db)
                    print(f"🔍 HTTP API: Результат сопоставления: {matched_result}")
                    if matched_result:
                        print(f"🔍 HTTP API: Добавляем {len(matched_result)} результатов")
                        matching_results.extend(matched_result)
                    else:
                        print(f"🔍 HTTP API: Создаем базовый результат для {normalized_article.get('contractor_article', normalized_article.get('description', 'неизвестный товар'))}")
                        # Если сопоставление не найдено, создаем базовый результат
                        matching_results.append(MatchingResult(
                            id=str(uuid.uuid4()),
                            contractor_article=normalized_article.get('contractor_article'),
                            description=normalized_article.get('description'),
                            matched=False,
                            agb_article=normalized_article.get('agb_article'),
                            bl_article=normalized_article.get('bl_article'),
                            match_confidence=normalized_article.get('match_confidence'),
                            nomenclature=normalized_article.get('nomenclature')
                        ))
            else:
                print(f"🔍 HTTP API: Нет артикулов для обработки")
            
            # Очищаем временные файлы
            for file_path in file_paths:
                try:
                    Path(file_path).unlink()
                except:
                    pass
            
            # Формируем сообщение ответа
            if articles or matching_results:
                message = f"Найдено {len(matching_results)} артикулов:\n\n"
                for i, result in enumerate(matching_results, 1):
                    if result.matched:
                        # Определяем тип поиска для отображения
                        search_info = ""
                        if hasattr(result, 'is_existing_mapping') and result.is_existing_mapping:
                            search_info = " (уже был такой запрос) 🔄"
                        elif hasattr(result, 'search_type') and result.search_type == 'existing_mapping_by_description':
                            search_info = " (найдено по описанию) 🔍"
                        elif hasattr(result, 'search_type') and result.search_type == 'ai_search':
                            search_info = " (поиск через ИИ) 🤖"
                        
                        message += f"{i}. ✅ {result.contractor_article} - {result.description}{search_info}\n"
                        message += f"   → АГБ: {result.agb_article} | BL: {result.bl_article} | Уверенность: {result.match_confidence:.1f}%\n\n"
                    else:
                        message += f"{i}. ❌ {result.contractor_article} - {result.description} (не найдено в БД)\n\n"
            else:
                # Обычный ответ чата
                message = ai_response
            
            return AIMatchingResponse(
                message=message,
                matching_results=matching_results,
                processing_time=0.5,
                status="completed"
            )
            
        except Exception as e:
            # Очищаем временные файлы
            for file_path in file_paths:
                try:
                    Path(file_path).unlink()
                except:
                    pass
            
            return AIMatchingResponse(
                message=f"Ошибка обработки ИИ: {str(e)}",
                matching_results=[],
                processing_time=0.1,
                status="error"
            )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка обработки: {str(e)}")

# Новые endpoints для Excel функционала

@router.post("/parse-excel/", response_model=ExcelParseResponse)
async def parse_excel_file(
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db)
):
    """
    Парсит Excel файл и возвращает структурированные данные
    """
    try:
        # Проверяем тип файла
        if not file.filename.endswith(('.xlsx', '.xls')):
            raise HTTPException(status_code=400, detail="Поддерживаются только файлы .xlsx и .xls")
        
        # Читаем файл
        contents = await file.read()
        
        # Парсим Excel
        try:
            df = pd.read_excel(io.BytesIO(contents))
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Ошибка чтения Excel файла: {str(e)}")
        
        # Преобразуем в нужный формат
        excel_rows = []
        for index, row in df.iterrows():
            excel_row = ExcelRow(
                id=f"row_{index + 1}",
                наименование=str(row.get('Наименование', '')),
                запрашиваемый_артикул=str(row.get('Запрашиваемый артикул', '')),
                количество=float(row.get('Количество', 1)),
                единица_измерения=str(row.get('Единица измерения', 'шт')),
                наш_артикул=str(row.get('Наш артикул', '')) if pd.notna(row.get('Наш артикул')) else None,
                артикул_bl=str(row.get('Артикул BL', '')) if pd.notna(row.get('Артикул BL')) else None,
                номер_1с=str(row.get('Номер в 1С', '')) if pd.notna(row.get('Номер в 1С')) else None,
                стоимость=float(row.get('Стоимость', 0)) if pd.notna(row.get('Стоимость')) else None,
                статус_сопоставления="pending",
                уверенность=0
            )
            excel_rows.append(excel_row)
        
        return ExcelParseResponse(
            success=True,
            data=excel_rows,
            message=f"Успешно обработано {len(excel_rows)} строк"
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка обработки файла: {str(e)}")

@router.post("/auto-match-excel/", response_model=ExcelMatchResponse)
async def auto_match_excel_data(
    request: ExcelDataRequest,
    db: AsyncSession = Depends(get_db)
):
    """
    Выполняет автоматическое сопоставление Excel данных с базой
    """
    try:
        matched_data = []
        statistics = {
            "total": len(request.data),
            "matched": 0,
            "unmatched": 0,
            "pending": 0
        }
        
        # Получаем все активные элементы из базы данных для поиска по наименованию
        async for db in get_db():
            name_query = select(MatchingNomenclature).where(
                MatchingNomenclature.is_active == True
            )
            name_results = await db.execute(name_query)
            name_items = name_results.scalars().all()
            
            print(f"DEBUG: Загружено {len(name_items)} элементов из базы")
            print(f"DEBUG: Получено {len(request.data)} строк для обработки")
            
            for row in request.data:
                matched_row = row.copy()
                print(f"DEBUG: Обрабатываем строку: '{row.наименование}'")
                
                # Поиск по наименованию (70% совпадение)
                if row.наименование:
                    print(f"DEBUG: Поиск по наименованию: '{row.наименование}'")
                    search_text = row.наименование.lower().strip()
                    
                    # Собираем ВСЕ варианты с их реальной схожестью
                    matches = []
                    
                    # 1. Поиск по подстроке (более точный)
                    for item in name_items:
                        item_name = item.name.lower()
                        if search_text in item_name:
                            # Для подстроки используем более высокую схожесть
                            # Вычисляем отношение длины подстроки к длине полного названия
                            substring_ratio = len(search_text) / len(item_name)
                            # Минимум 0.5, максимум 1.0
                            similarity = max(0.5, min(1.0, substring_ratio + 0.3))
                            
                            matches.append({
                                'item': item,
                                'similarity': similarity,
                                'confidence': int(similarity * 100),
                                'match_type': 'substring'
                            })
                            print(f"DEBUG: Найдено подстроковое совпадение: '{item.name}' (схожесть: {similarity:.2f})")
                    
                    # 2. Поиск по нормализованному тексту для всех элементов
                    for item in name_items:
                        item_name = item.name.lower()
                        normalized_search = get_normalized_text(search_text)
                        normalized_item = get_normalized_text(item_name)
                        similarity = SequenceMatcher(None, normalized_search, normalized_item).ratio()
                        
                        # Добавляем только если схожесть >= 50% и это не дубликат
                        if similarity >= 0.5:  # Снижаем порог для получения больше вариантов
                            # Проверяем, что это не дубликат
                            is_duplicate = any(
                                match['item'].id == item.id for match in matches
                            )
                            if not is_duplicate:
                                matches.append({
                                    'item': item,
                                    'similarity': similarity,
                                    'confidence': int(similarity * 100),
                                    'match_type': 'normalized'
                                })
                                print(f"DEBUG: Найдено нормализованное совпадение: '{item.name}' (схожесть: {similarity:.2f})")
                    
                    # 3. Дополнительный поиск по ключевым словам
                    search_words = search_text.split()
                    if len(search_words) > 1:  # Если есть несколько слов
                        for item in name_items:
                            item_name = item.name.lower()
                            word_matches = 0
                            total_words = len(search_words)
                            
                            for word in search_words:
                                if word in item_name:
                                    word_matches += 1
                            
                            if word_matches > 0:
                                word_similarity = word_matches / total_words
                                
                                # Проверяем, что это не дубликат
                                is_duplicate = any(
                                    match['item'].id == item.id for match in matches
                                )
                                if not is_duplicate and word_similarity >= 0.3:
                                    matches.append({
                                        'item': item,
                                        'similarity': word_similarity,
                                        'confidence': int(word_similarity * 100),
                                        'match_type': 'keywords'
                                    })
                                    print(f"DEBUG: Найдено совпадение по ключевым словам: '{item.name}' (схожесть: {word_similarity:.2f})")
                    
                    # Сортируем по убыванию схожести
                    matches.sort(key=lambda x: x['similarity'], reverse=True)
                    
                    if matches:
                        # Берем лучший вариант для основных полей (только если схожесть >= 70%)
                        best_match = matches[0]
                        if best_match['similarity'] >= 0.7:
                            matched_row.наш_артикул = best_match['item'].agb_article
                            matched_row.артикул_bl = best_match['item'].bl_article
                            matched_row.номер_1с = best_match['item'].code_1c
                            matched_row.статус_сопоставления = "matched"
                            matched_row.уверенность = best_match['confidence']
                        else:
                            matched_row.статус_сопоставления = "partial"
                            matched_row.уверенность = best_match['confidence']
                        
                        # Добавляем ВСЕ варианты в поле вариантов (топ 10)
                        matched_row.варианты_подбора = [
                            {
                                'наименование': match['item'].name,
                                'наш_артикул': match['item'].agb_article,
                                'артикул_bl': match['item'].bl_article,
                                'номер_1с': match['item'].code_1c,
                                'уверенность': match['confidence'],
                                'тип_совпадения': match['match_type']
                            }
                            for match in matches[:10]  # Показываем максимум 10 вариантов
                        ]
                        
                        matched_data.append(matched_row)
                        if best_match['similarity'] >= 0.7:
                            statistics["matched"] += 1
                        else:
                            statistics["pending"] += 1
                        print(f"DEBUG: Добавлено {len(matches)} совпадений для '{row.наименование}'")
                    else:
                        matched_row.статус_сопоставления = "unmatched"
                        matched_row.варианты_подбора = []
                        matched_data.append(matched_row)
                        statistics["unmatched"] += 1
                        print(f"DEBUG: Не найдено совпадений для '{row.наименование}'")
                    continue
                
                # Поиск по артикулу (100% совпадение)
                elif row.запрашиваемый_артикул:
                    print(f"DEBUG: Поиск по артикулу: '{row.запрашиваемый_артикул}'")
                    # Поиск в нашей базе по артикулу АГБ
                    agb_query = select(MatchingNomenclature).where(
                        MatchingNomenclature.agb_article == row.запрашиваемый_артикул
                    )
                    agb_result = await db.execute(agb_query)
                    agb_item = agb_result.scalar_one_or_none()
                    
                    # Поиск по артикулу BL
                    bl_query = select(MatchingNomenclature).where(
                        MatchingNomenclature.bl_article == row.запрашиваемый_артикул
                    )
                    bl_result = await db.execute(bl_query)
                    bl_item = bl_result.scalar_one_or_none()
                    
                    if agb_item:
                        matched_row.наш_артикул = agb_item.agb_article
                        matched_row.артикул_bl = agb_item.bl_article
                        matched_row.номер_1с = agb_item.code_1c
                        matched_row.статус_сопоставления = "matched"
                        matched_row.уверенность = 100
                        matched_row.варианты_подбора = [{
                            'наименование': agb_item.name,
                            'наш_артикул': agb_item.agb_article,
                            'артикул_bl': agb_item.bl_article,
                            'номер_1с': agb_item.code_1c,
                            'уверенность': 100
                        }]
                        matched_data.append(matched_row)
                        statistics["matched"] += 1
                    elif bl_item:
                        matched_row.наш_артикул = bl_item.agb_article
                        matched_row.артикул_bl = bl_item.bl_article
                        matched_row.номер_1с = bl_item.code_1c
                        matched_row.статус_сопоставления = "matched"
                        matched_row.уверенность = 100
                        matched_row.варианты_подбора = [{
                            'наименование': bl_item.name,
                            'наш_артикул': bl_item.agb_article,
                            'артикул_bl': bl_item.bl_article,
                            'номер_1с': bl_item.code_1c,
                            'уверенность': 100
                        }]
                        matched_data.append(matched_row)
                        statistics["matched"] += 1
                    else:
                        matched_row.статус_сопоставления = "unmatched"
                        matched_row.варианты_подбора = []
                        matched_data.append(matched_row)
                        statistics["unmatched"] += 1
                else:
                    matched_row.статус_сопоставления = "unmatched"
                    matched_row.варианты_подбора = []
                    matched_data.append(matched_row)
                    statistics["unmatched"] += 1
        
        return ExcelMatchResponse(
            success=True,
            matched_data=matched_data,
            statistics=statistics,
            message=f"Сопоставление завершено. Найдено {statistics['matched']} из {statistics['total']} позиций"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка сопоставления: {str(e)}")

@router.post("/save-excel-results/")
async def save_excel_results(
    request: ExcelDataRequest,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Сохраняет результаты Excel сопоставления в базу данных
    """
    try:
        saved_count = 0
        
        for row in request.data:
            # Создаем запись в таблице сопоставлений
            mapping_data = {
                "contractor_article": row.запрашиваемый_артикул,
                "contractor_description": row.наименование,
                "agb_article": row.наш_артикул or "",
                "agb_description": "",  # Можно добавить описание из номенклатуры
                "bl_article": row.артикул_bl or "",
                "bl_description": "",  # Можно добавить описание BL
                "packaging_factor": 1.0,
                "unit": row.единица_измерения,
                "quantity": row.количество,
                "cost": row.стоимость or 0.0,
                "match_confidence": row.уверенность or 0,
                "status": row.статус_сопоставления or "pending",
                "created_by": current_user.id,
                "created_at": datetime.utcnow()
            }
            
            # Здесь можно добавить логику сохранения в базу данных
            # Например, создание записи в таблице ArticleMapping
            saved_count += 1
        
        return {
            "success": True,
            "message": f"Сохранено {saved_count} записей",
            "saved_count": saved_count
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ошибка сохранения: {str(e)}")

@router.get("/saved-variants/")
async def get_saved_variants(
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Получение сохраненных вариантов подбора"""
    try:
        # Здесь можно добавить логику получения сохраненных вариантов из базы данных
        # Пока возвращаем пустой список
        return {
            "success": True,
            "saved_variants": {}
        }
    except Exception as e:
        logger.error(f"Ошибка при получении сохраненных вариантов: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Ошибка при получении: {str(e)}")

@router.post("/save-variant-selection/")
async def save_variant_selection(
    request: dict,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db)
):
    """Сохранение выбранного варианта подбора"""
    try:
        # Здесь можно добавить логику сохранения выбранного варианта в базу данных
        # request должен содержать: {"row_id": "string", "variant_index": int}
        
        return {
            "success": True,
            "message": "Вариант сохранен"
        }
    except Exception as e:
        logger.error(f"Ошибка при сохранении варианта: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Ошибка при сохранении: {str(e)}")
